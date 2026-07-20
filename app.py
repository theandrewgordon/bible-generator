# test08-28-2025
from flask import Flask, Response, render_template, request, send_file, send_from_directory, redirect, url_for, session, flash, jsonify, g, after_this_request, has_request_context
from flask_dance.contrib.google import make_google_blueprint, google
from datetime import datetime

import os, json, re, traceback
import logging
import sys
import uuid
import socket
import ipaddress
from html.parser import HTMLParser
from urllib.parse import urlparse, urljoin, urlunparse
from urllib.request import Request, urlopen, build_opener, HTTPRedirectHandler
from zipfile import ZipFile
import threading
from io import BytesIO
from datetime import datetime, timedelta, timezone
from firebase_admin import firestore
from werkzeug.utils import secure_filename
from werkzeug.middleware.proxy_fix import ProxyFix
import time
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from verse_helpers import (
    request_verse_data,
    parse_and_clean_json,
    save_json_to_file,
    ai_validate_custom_text,
    normalize_reference_title,
    normalize_verse_data,
    preserve_letter_suffix,
)
from build_pdf import generate_pdf
from PIL import Image, ImageDraw, ImageFont
try:
    import markdown2  # type: ignore
except Exception:
    markdown2 = None

# Extracted services/utilities
from faithsparks.services.firestore import (
    db,
    firebase_init_diagnostic,
    init_firebase,
    validate_firebase_credentials,
)
from faithsparks.services.storage import upload_to_storage, signed_url_for_path
from faithsparks.services.collections import get_collections, get_collection_meta, get_collection_verses, COLLECTIONS
from faithsparks.services.usage import _month_key, _get_user_plan, _get_usage, _quota_for_plan, _update_usage, _get_free_slugs
from faithsparks.services.users import get_user_doc
from faithsparks.services.themes import THEMES, get_theme_vars, list_all_themes, get_theme_selection
from faithsparks.services.stripe_svc import (
    stripe, STRIPE_SECRET_KEY, STRIPE_PUBLISHABLE_KEY,
    STRIPE_PRICE_FAMILY, STRIPE_PRICE_CLASSROOM,
    STRIPE_PRICE_FAMILY_MONTHLY, STRIPE_PRICE_FAMILY_ANNUAL,
    STRIPE_PRICE_CLASSROOM_MONTHLY, STRIPE_PRICE_CLASSROOM_ANNUAL,
    STRIPE_WEBHOOK_SECRET, resolve_price_id as _resolve_price_id
)
from faithsparks.util.slug import normalize_slug
from faithsparks.services import analytics as analytics_svc
from faithsparks.services.rate_limit import check_rate_limit
from faithsparks.util.request_utils import get_client_ip, get_request_payload, log_request_summary
from faithsparks.views.worksheets import MAX_WORKSHEETS_PER_REQUEST

# Map passwords to pack metadata
PACKS = {
    # ESV – Print
    "sparks-esv-print": {
        "id": "esv_print",
        "name": "ESV Print Handwriting (30 Worksheets)",
        "filename": "30_Pack_ESV.pdf",
    },
    # ESV – Cursive
    "sparks-esv-cursive": {
        "id": "esv_cursive",
        "name": "ESV Cursive Handwriting (30 Worksheets)",
        "filename": "30_Pack_ESV_Cursive.pdf",
    },
    # KJV – Print
    "sparks-kjv-print": {
        "id": "kjv_print",
        "name": "KJV Print Handwriting (30 Worksheets)",
        "filename": "30_Pack_KJV.pdf",
    },
    # KJV – Cursive
    "sparks-kjv-cursive": {
        "id": "kjv_cursive",
        "name": "KJV Cursive Handwriting (30 Worksheets)",
        "filename": "30_Pack_KJV_Cursive.pdf",
    },
    # Mega bundle (all 4)
    "sparks-mega-bundle": {
        "id": "mega_bundle",
        "name": "Mega Bundle – All 4 Packs (120 Worksheets)",
        "filename": "30_Pack_MegaPack.zip",
    },
}

_CONFIG_CACHE: dict[str, dict] = {}
_CONFIG_TTL_SECS = 60
CLEANUP_INTERVAL_S = 60 * 60
CLEANUP_MAX_AGE_S = 7 * 24 * 60 * 60
_LAST_CLEANUP: float = 0.0

def _get_cached_config(doc_id: str) -> dict | None:
    """Cache repeated config reads for a short TTL."""
    now = time.time()
    entry = _CONFIG_CACHE.get(doc_id)
    if entry and now - entry.get("ts", 0) < _CONFIG_TTL_SECS:
        return entry["data"]
    if not db:
        return {}
    try:
        snap = db.collection("config").document(doc_id).get()
        data = snap.to_dict() if snap.exists else {}
    except Exception:
        data = {}
    _CONFIG_CACHE[doc_id] = {"data": data, "ts": now}
    return data

def _refresh_owned_packs(email: str | None) -> set[str]:
    """Cache the list of user-owned packs in session to avoid repeated reads."""
    if not (db and email):
        session.pop("user_owned_packs", None)
        return set()
    try:
        docs = db.collection("users").document(email).collection("purchases").stream()
        packs = {doc.id for doc in docs}
    except Exception:
        packs = set()
    session["user_owned_packs"] = list(packs)
    return packs

def _should_fetch_usage(path: str) -> bool:
    """Only refresh usage info for high-traffic pages."""
    if not path:
        return False
    for prefix in ("/generate", "/browse", "/prints", "/history", "/plus", "/games"):
        if path.startswith(prefix):
            return True
    return False

def _cleanup_output_dirs():
    global _LAST_CLEANUP
    now = time.time()
    if now - _LAST_CLEANUP < CLEANUP_INTERVAL_S:
        return
    _LAST_CLEANUP = now
    dirs = ["output", "output/thumbs", "output/packs"]
    cutoff = now - CLEANUP_MAX_AGE_S
    for base in dirs:
        path = pathlib.Path(base)
        if not path.exists():
            continue
        for child in path.iterdir():
            try:
                if child.is_dir():
                    continue
                if child.stat().st_mtime < cutoff:
                    child.unlink()
            except Exception:
                continue
# --- App Setup ---
# --- Environment / config flags (MUST be defined before use) ---
APP_ENV = os.getenv("APP_ENV", "dev").lower()
PRIMARY_DOMAIN = os.getenv("PRIMARY_DOMAIN", "faithsparksprintables.com")

def _is_local_storage_allowed() -> bool:
    """Return True if local storage fallback is allowed (development-only by default)."""
    if os.getenv("USE_LOCAL_STORAGE", "").lower() in {"1", "true", "yes"}:
        return True
    return APP_ENV not in {"prod", "production"}

def _compute_static_version() -> str:
    """Cache-bust token for static assets. Changes whenever the deploy (or the
    static files) change, so returning visitors don't get stale CSS/JS."""
    explicit = os.getenv("RENDER_GIT_COMMIT") or os.getenv("STATIC_VERSION")
    if explicit:
        return explicit[:12]
    latest = 0.0
    for rel in (
        "static/theme.css",
        "static/darkmode.js",
        "static/admin.css",
        "static/bible_bee.css",
        "static/bible_bee.js",
    ):
        try:
            latest = max(latest, os.path.getmtime(rel))
        except OSError:
            pass
    return str(int(latest)) if latest else "1"

STATIC_VERSION = _compute_static_version()

# How long the navbar may serve is_pro/plan from the session cache before
# re-reading Firestore. Short enough that plan changes show up quickly.
_USER_FLAGS_TTL = 120.0  # seconds

# --- App Setup ---
app = Flask(__name__)

# Fail fast in production if secret or required cloud storage config is missing
if APP_ENV in {"prod", "production"}:
    if not os.getenv("FLASK_SECRET_KEY"):
        raise RuntimeError("FLASK_SECRET_KEY must be set in production")
    if not os.getenv("FIREBASE_CREDS_JSON"):
        raise RuntimeError("FIREBASE_CREDS_JSON must be set in production to configure Firestore")
    try:
        validate_firebase_credentials()
    except Exception as exc:
        raise RuntimeError("FIREBASE_CREDS_JSON is invalid") from exc

app.secret_key = os.getenv("FLASK_SECRET_KEY", "dev-only-secret")


# Structured logging to stdout for easier aggregation
handler = logging.StreamHandler(sys.stdout)
handler.setFormatter(logging.Formatter("%(asctime)s %(levelname)s %(message)s"))
if not app.logger.handlers:
    app.logger.addHandler(handler)
app.logger.setLevel(logging.INFO)
app.logger.propagate = False
if APP_ENV in {"prod", "production"}:
    app.logger.info("Firestore credentials validated; client initialization deferred to worker")


# Always prefer https URLs when generating links
app.config.update(PREFERRED_URL_SCHEME="https")

# Only pin cookies to the apex domain in production
if APP_ENV in {"prod", "production"}:
    app.config["SESSION_COOKIE_DOMAIN"] = f".{PRIMARY_DOMAIN}"

# Respect proxy headers (Render/Cloudflare) when enabled
enable_proxy_fix = os.getenv(
    "ENABLE_PROXY_FIX",
    "1" if APP_ENV in {"prod", "production"} else "0",
).lower() in {"1", "true", "yes"}
if enable_proxy_fix:
    app.wsgi_app = ProxyFix(app.wsgi_app, x_for=1, x_proto=1, x_host=1, x_prefix=1)

@app.route('/manifest.webmanifest')
def pwa_manifest():
    """Serve the PWA manifest with minimal caching for quick updates."""
    response = send_from_directory('static', 'manifest.webmanifest')
    response.headers['Cache-Control'] = 'no-cache'
    response.headers['Content-Type'] = 'application/manifest+json'
    return response

@app.route('/service-worker.js')
def service_worker():
    """Serve the service worker from the app root for full-scope control."""
    response = send_from_directory('static', 'service-worker.js')
    response.headers['Cache-Control'] = 'no-cache'
    return response


@app.route('/robots.txt')
def robots_txt():
    body = "\n".join(
        [
            "User-agent: *",
            "Allow: /",
            "Sitemap: https://faithsparksprintables.com/sitemap.xml",
            "",
        ]
    )
    return Response(body, mimetype="text/plain")

@app.route("/speeddie")
@app.route("/speeddie/")
def speeddie_app():
    return send_from_directory("speeddie", "index.html")


@app.route("/speeddie/<path:filename>")
def speeddie_static(filename):
    return send_from_directory("speeddie", filename)

@app.route("/verse-of-the-week")
def verse_of_the_week():
    """Weekly memory verse + one-click worksheet — a reason to return each week."""
    from faithsparks.util.verse_of_week import get_verse_of_week
    votw = get_verse_of_week()
    preview = None
    try:
        from faithsparks.services.scripture import fetch_verse_text
        preview = fetch_verse_text(votw["reference"], "web")  # public-domain, safe to display
    except Exception:
        preview = None
    return render_template("verse_of_week.html", votw=votw, preview=preview)


# Jinja filter for Markdown
def _md(text: str) -> str:
    try:
        if not text:
            return ""
        if markdown2:
            return markdown2.markdown(text)
        return text
    except Exception:
        return text

app.jinja_env.filters["markdown"] = _md

# Recommended cookie settings (don't break localhost/dev)
app.config["SESSION_COOKIE_SECURE"] = (APP_ENV in {"prod", "production"})
app.config["SESSION_COOKIE_SAMESITE"] = "Lax"


# Helpful OAuth env flags (no-op if already set)
os.environ.setdefault("OAUTHLIB_RELAX_TOKEN_SCOPE", "1")
# Do NOT set OAUTHLIB_INSECURE_TRANSPORT to 1 in production

# --- Google Auth ---
google_bp = make_google_blueprint(
    client_id=os.environ.get("GOOGLE_OAUTH_CLIENT_ID"),
    client_secret=os.environ.get("GOOGLE_OAUTH_CLIENT_SECRET"),
    redirect_to="oauth_finish",                    # <— was "index"
    scope=["openid", "https://www.googleapis.com/auth/userinfo.email", "https://www.googleapis.com/auth/userinfo.profile"]
)
app.register_blueprint(google_bp, url_prefix="/login")

@app.before_request
def add_request_id():
    if not getattr(g, "req_id", None):
        g.req_id = request.headers.get("X-Request-ID") or str(uuid.uuid4())

@app.before_request
def cleanup_old_outputs():
    _cleanup_output_dirs()

@app.before_request
def force_primary_domain():
    if APP_ENV not in {"prod", "production"}:
        return
    host = request.host.split(":")[0]
    if host == PRIMARY_DOMAIN:
        return
    parsed = urlparse(request.url)
    target = parsed._replace(scheme="https", netloc=PRIMARY_DOMAIN)
    return redirect(urlunparse(target), code=301)

def is_safe_url(target: str) -> bool:
    if not target:
        return False
    ref = urlparse(request.host_url)
    test = urlparse(urljoin(request.host_url, target))
    return (test.scheme in ("http", "https")) and (ref.netloc == test.netloc)

@app.before_request
def load_user_info():
    # When Flask-Dance blueprint handles /login/google*, capture ?next=...
    if request.blueprint == "google" or request.path.startswith("/login/google"):
        nxt = request.args.get("next")
        if nxt and is_safe_url(nxt):
            session["post_login_next"] = nxt
        return  # don't do user_info fetch mid-dance

    if google.authorized:
        if "user_info" not in session:
            try:
                resp = google.get("/oauth2/v1/userinfo")
                if resp.ok:
                    session["user_info"] = resp.json()
                    session["user_email"] = session["user_info"].get("email")
                    session["clear_storage"] = True
                    _refresh_owned_packs(session["user_email"])
            except Exception:
                # Token expired or revoked — clear session so user is prompted to re-auth
                session.pop("user_info", None)
                session.pop("user_email", None)
                session.pop("user_owned_packs", None)
        elif not session.get("user_owned_packs"):
            _refresh_owned_packs(session.get("user_email"))
    else:
        session.pop("user_info", None)
        session.pop("user_email", None)
        session.pop("user_owned_packs", None)
        session.pop("_uc", None)


# -----------------------------
# CSRF (lightweight, no Flask-WTF)
# -----------------------------
import secrets
from flask import abort

CSRF_SESSION_KEY = "_csrf_token"
CSRF_HEADER_NAME = "X-CSRF-Token"
CSRF_ALT_HEADER_NAME = "X-CSRFToken"
CSRF_FORM_FIELD = "csrf_token"

def _get_csrf_token() -> str:
    """Return (and create if missing) a session-bound CSRF token."""
    tok = session.get(CSRF_SESSION_KEY)
    if not tok:
        tok = secrets.token_urlsafe(32)
        session[CSRF_SESSION_KEY] = tok
    return tok

def _constant_time_eq(a: str, b: str) -> bool:
    try:
        return secrets.compare_digest(a or "", b or "")
    except Exception:
        return False

# Routes/endpoints to skip CSRF checks (webhooks + OAuth + static)
_CSRF_EXEMPT_PREFIXES = (
    "/stripe/webhook",
    "/login/google",     # flask-dance endpoints under this prefix
    "/oauth",            # your /oauth/finish etc
    "/static/",
    "/service-worker.js",
    "/manifest.webmanifest",
)

@app.before_request
def csrf_protect():
    """Validate CSRF token on mutating requests."""
    if request.method in ("GET", "HEAD", "OPTIONS", "TRACE"):
        return

    path = request.path or ""
    for p in _CSRF_EXEMPT_PREFIXES:
        if path.startswith(p):
            return

    # Accept token from header (AJAX/fetch) or form field
    sent = (
        request.headers.get(CSRF_HEADER_NAME)
        or request.headers.get(CSRF_ALT_HEADER_NAME)
        or request.form.get(CSRF_FORM_FIELD)
    )

    expected = session.get(CSRF_SESSION_KEY)

    if not expected:
        # session missing token -> force creation and fail this request
        _get_csrf_token()
        abort(403)

    if not _constant_time_eq(sent, expected):
        abort(403)

@app.context_processor
def inject_csrf():
    """Expose csrf_token() helper to all templates."""
    return {"csrf_token": _get_csrf_token}


@app.before_request
def track_visit():
    if request.method not in ("GET", "HEAD"):
        return
    endpoint = request.endpoint or ""
    path = request.path or ""
    if endpoint == "static" or path.startswith("/static/"):
        return
    try:
        ip = get_client_ip()
        ua = request.headers.get("User-Agent", "")
        analytics_svc.record_visit(ip, ua)
    except Exception:
        app.logger.debug("Failed to record visit", exc_info=True)


@app.after_request
def add_correlation_headers(resp):
    req_id = getattr(g, "req_id", None)
    if req_id:
        resp.headers["X-Request-ID"] = req_id
    return resp


app.teardown_appcontext(analytics_svc.close_db)

# quotas and usage moved to yourapp.services.usage

# --- Helpers ---
def login_required(func):
    from functools import wraps
    @wraps(func)
    def wrapper(*args, **kwargs):
        if not google.authorized:
            # keep path + query (e.g., /generate?v=John%203:16%20(ESV))
            next_url = request.full_path.rstrip("?") if request.query_string else request.path
            return redirect(url_for("google.login", next=next_url))
        return func(*args, **kwargs)
    return wrapper

def admin_required(func):
    from functools import wraps
    @wraps(func)
    def wrapper(*args, **kwargs):
        if not google.authorized:
            return redirect(url_for("google.login"))
        email = session.get('user_email')
        if not is_admin_email(email):
            return "Forbidden", 403
        return func(*args, **kwargs)
    return wrapper

def is_public_browse_enabled() -> bool:
    return os.getenv('PUBLIC_BROWSE', '0') in ('1','true','True','yes','on')

def _fmt_dt(ts):
    try:
        return ts.strftime('%Y-%m-%d %H:%M') if ts else None
    except Exception:
        try:
            # Fallback to string
            return str(ts)
        except Exception:
            return None


def get_pack_by_password(raw_password: str):
    """Normalize and resolve a pack from a user-entered password."""
    if not raw_password:
        return None
    key = raw_password.strip().lower()
    return PACKS.get(key)


def get_pack_by_id(pack_id: str):
    """Return pack metadata by its id."""
    for data in PACKS.values():
        if data.get("id") == pack_id:
            return data
    return None


def _user_has_pack(user_email: str | None, pack_id: str) -> bool:
    if not (user_email and pack_id):
        return False
    owned = set(session.get("user_owned_packs") or [])
    if pack_id in owned:
        return True
    if not db:
        return False
    try:
        doc = (
            db.collection("users")
            .document(user_email)
            .collection("purchases")
            .document(pack_id)
            .get()
        )
        exists = doc.exists
        if exists:
            owned.add(pack_id)
            session["user_owned_packs"] = list(owned)
        return exists
    except Exception:
        return False


# --- Downloads portal ---
@app.route("/downloads", methods=["GET", "POST"])
def downloads():
    """
    Password-gated portal for pre-made worksheet packs.
    """
    error = None
    pack = None

    if request.method == "POST":
        password = request.form.get("password", "")
        pack = get_pack_by_password(password)
        if not pack:
            error = "That password doesn’t match any product. Please double-check and try again."
        else:
            session["unlocked_pack_id"] = pack["id"]
            flash(f"Unlocked {pack['name']}.", "success")

    if not pack:
        unlocked = session.get("unlocked_pack_id")
        if unlocked:
            pack = get_pack_by_id(unlocked)

    user_email = session.get("user_email")

    return render_template("downloads.html", error=error, pack=pack, user_email=user_email)


@app.route("/downloads/file/<pack_id>")
def download_file_pack(pack_id):
    """
    Serve the actual bundle for the unlocked pack.
    """
    unlocked_id = session.get("unlocked_pack_id")
    user_email = session.get("user_email")
    if not ((unlocked_id and unlocked_id == pack_id) or _user_has_pack(user_email, pack_id)):
        flash("Please enter your product password first.")
        return redirect(url_for("downloads"))
    if not unlocked_id and _user_has_pack(user_email, pack_id):
        session["unlocked_pack_id"] = pack_id

    pack = get_pack_by_id(pack_id)
    if not pack:
        flash("We couldn’t find that product.")
        return redirect(url_for("downloads"))

    bundles_dir = os.path.join(app.root_path, "static", "bundles")
    filename = pack["filename"]
    file_path = os.path.join(bundles_dir, filename)

    if not os.path.exists(file_path):
        app.logger.warning("Bundle missing on disk: %s", file_path)
        flash("We couldn’t find your download file. Please contact support.")
        return redirect(url_for("downloads"))

    return send_from_directory(
        bundles_dir,
        filename,
        as_attachment=True,
        download_name=filename,
    )


@app.route("/downloads/claim", methods=["POST"])
@login_required
def claim_pack():
    """
    Save the unlocked pack to the signed-in user's library.
    """
    user_email = session.get("user_email")
    pack_id = request.form.get("pack_id")
    if not pack_id:
        flash("We couldn’t tell which pack to save.")
        return redirect(url_for("downloads"))

    pack_meta = get_pack_by_id(pack_id)
    if not pack_meta:
        flash("That pack doesn’t exist.")
        return redirect(url_for("downloads"))

    try:
        doc_ref = (
            db.collection("users")
            .document(user_email)
            .collection("purchases")
            .document(pack_id)
        )
        doc_ref.set(
            {
                "pack_id": pack_id,
                "name": pack_meta["name"],
                "filename": pack_meta["filename"],
                "created_at": firestore.SERVER_TIMESTAMP,
                "source": "downloads_portal",
            },
            merge=True,
        )
        flash("Pack saved to your Faith Sparks Library! You can access it anytime from your account.")
        owned = set(session.get("user_owned_packs") or [])
        owned.add(pack_id)
        session["user_owned_packs"] = list(owned)
    except Exception as e:
        app.logger.exception("Failed to save pack to library: %s", e)
        flash("We couldn’t save this to your library. Please try again.")

    return redirect(url_for("downloads"))


@app.route("/downloads/restore/<pack_id>")
@login_required
def restore_pack(pack_id):
    """
    Restore an owned pack from Library, set session unlock, and bounce to portal.
    """
    user_email = session.get("user_email")
    if not _user_has_pack(user_email, pack_id):
        flash("That pack is not in your library yet.")
        return redirect(url_for("prints"))
    session["unlocked_pack_id"] = pack_id
    flash("Pack ready to download.")
    return redirect(url_for("downloads"))

# normalize_slug moved to yourapp.util.slug

def extract_version_from_text(text, fallback_version):
    norm_fallback = (fallback_version or "esv").lower().strip()
    fallback_version = "esv" if norm_fallback in ("", "auto") else norm_fallback
    match = re.search(r'\(([A-Za-z0-9]{2,12})\)\s*$', text.strip())
    if match:
        version = match.group(1).lower()
        verse = text[:match.start()].strip()
    else:
        version = fallback_version
        verse = text.strip()
    return version, normalize_reference_title(verse)


def _boolish(value, default=False):
    if isinstance(value, bool):
        return value
    if isinstance(value, str):
        return value.lower().strip() in {"1", "true", "yes", "on"}
    return bool(value) if value is not None else default


SONGS_DIR = os.path.join(app.root_path, "songs")  # kept for local fallback

# --- Worship song persistence helpers ---
_WORSHIP_COLLECTION = "worship_songs"
_WORSHIP_SETLIST_COLLECTION = "worship_setlists"
_WORSHIP_SCOPE_COLLECTION = "worship_scopes"
_WORSHIP_CHURCH_COLLECTION = "worship_churches"
_DEFAULT_WORSHIP_SCOPE = "default"
_worship_seeded = False
# Per-scope cache: {scope_id: {"data": list[dict], "ts": float}}
_worship_songs_cache: dict[str, dict] = {}
_WORSHIP_CACHE_TTL = 60.0  # seconds


def _current_worship_scope() -> str:
    raw = None
    if has_request_context():
        raw = session.get("worship_church_id") or session.get("worship_scope")
        if not raw and db and session.get("user_email"):
            try:
                user_doc = db.collection("users").document(session["user_email"]).get()
                data = user_doc.to_dict() if user_doc.exists else {}
                candidate = data.get("worshipChurchId")
                if candidate and _user_can_access_worship_church(candidate):
                    raw = candidate
                    session["worship_church_id"] = candidate
            except Exception:
                raw = None
    raw = raw or os.getenv("WORSHIP_SCOPE_ID") or _DEFAULT_WORSHIP_SCOPE
    scope = _slugify_worship_token(str(raw))
    return scope or _DEFAULT_WORSHIP_SCOPE


def _normalize_worship_invite_code(value: str) -> str:
    return re.sub(r"[^A-Z0-9]+", "", str(value or "").upper())[:24]


def _make_worship_invite_code(name: str) -> str:
    prefix = re.sub(r"[^A-Z0-9]+", "", str(name or "").upper())[:4] or "FS"
    return f"{prefix}{secrets.token_hex(3).upper()}"


def _user_can_access_worship_church(church_id: str) -> bool:
    if not (db and church_id and session.get("user_email")):
        return False
    try:
        member = (
            db.collection(_WORSHIP_CHURCH_COLLECTION)
            .document(_slugify_worship_token(church_id))
            .collection("members")
            .document(session["user_email"])
            .get()
        )
        return member.exists
    except Exception:
        return False


def _load_worship_churches() -> list[dict]:
    if not (db and session.get("user_email")):
        return []
    email = session["user_email"]
    memberships: list[dict] = []
    try:
        user_doc = db.collection("users").document(email).get()
        user_data = user_doc.to_dict() if user_doc.exists else {}
        church_ids = list(dict.fromkeys(user_data.get("worshipChurchIds") or []))
        if not church_ids:
            try:
                member_docs = db.collection_group("members").where(filter=firestore.FieldFilter("email", "==", email)).stream()
                for member_doc in member_docs:
                    church_ref = member_doc.reference.parent.parent
                    if church_ref:
                        church_ids.append(church_ref.id)
            except Exception:
                church_ids = []
            church_ids = list(dict.fromkeys(church_ids))
        for church_id in church_ids:
            church_ref = db.collection(_WORSHIP_CHURCH_COLLECTION).document(_slugify_worship_token(church_id))
            church_doc = church_ref.get()
            if not church_doc.exists:
                continue
            member_doc = church_ref.collection("members").document(email).get()
            if not member_doc.exists:
                continue
            church = church_doc.to_dict() or {}
            church_id = church.get("id") or church_doc.id
            memberships.append(
                {
                    "id": church_id,
                    "name": church.get("name") or church_id.replace("-", " ").title(),
                    "invite_code": church.get("invite_code") or church.get("inviteCode") or "",
                    "role": (member_doc.to_dict() or {}).get("role") or "member",
                }
            )
    except Exception as exc:
        app.logger.warning("_load_worship_churches error: %s", exc)
    memberships.sort(key=lambda item: item["name"].lower())
    return memberships


def _current_worship_church_context() -> dict:
    churches = _load_worship_churches()
    current_id = _current_worship_scope()
    current = next((church for church in churches if church["id"] == current_id), None)
    if not current:
        current = {"id": _DEFAULT_WORSHIP_SCOPE, "name": "Shared Library", "invite_code": "", "role": "member"}
    return {"current": current, "churches": churches}


def _set_current_worship_church(church_id: str) -> None:
    church_id = _slugify_worship_token(church_id)
    session["worship_church_id"] = church_id
    if db and session.get("user_email"):
        update_data = {"worshipChurchId": church_id}
        if church_id != _DEFAULT_WORSHIP_SCOPE:
            update_data["worshipChurchIds"] = firestore.ArrayUnion([church_id])
        db.collection("users").document(session["user_email"]).set(update_data, merge=True)


def _worship_songs_ref(client=None):
    if client is None:
        client = db
    return client.collection(_WORSHIP_SCOPE_COLLECTION).document(_current_worship_scope()).collection(_WORSHIP_COLLECTION)


def _legacy_worship_songs_ref(client=None):
    if client is None:
        client = db
    return client.collection(_WORSHIP_COLLECTION)


def _worship_setlists_ref():
    return db.collection(_WORSHIP_SCOPE_COLLECTION).document(_current_worship_scope()).collection(_WORSHIP_SETLIST_COLLECTION)


def _legacy_worship_setlists_ref():
    return db.collection(_WORSHIP_SETLIST_COLLECTION)


def _worship_song_refs_for_read(client=None):
    refs = [_worship_songs_ref(client)]
    if _current_worship_scope() == _DEFAULT_WORSHIP_SCOPE:
        refs.insert(0, _legacy_worship_songs_ref(client))
    return refs


def _worship_setlist_refs_for_read():
    refs = [_worship_setlists_ref()]
    if _current_worship_scope() == _DEFAULT_WORSHIP_SCOPE:
        refs.insert(0, _legacy_worship_setlists_ref())
    return refs


def _seed_worship_from_files() -> None:
    """One-time per-process seed: load /songs/*.json into Firestore if collection is empty."""
    global _worship_seeded
    if _worship_seeded:
        return
    _worship_seeded = True
    if not db:
        return
    try:
        if list(_worship_songs_ref().limit(1).stream()) or (
            _current_worship_scope() == _DEFAULT_WORSHIP_SCOPE and list(_legacy_worship_songs_ref().limit(1).stream())
        ):
            return  # already has documents
        songs_folder = Path(app.root_path) / "songs"
        if not songs_folder.is_dir():
            return
        for fp in sorted(songs_folder.glob("*.json")):
            try:
                with open(fp, "r", encoding="utf-8") as f:
                    data = json.load(f)
                song_id = data.get("id") or fp.stem
                data["id"] = song_id
                _worship_songs_ref().document(song_id).set(data)
                app.logger.info("Seeded worship song: %s", song_id)
            except Exception as exc:
                app.logger.warning("Worship seed skip %s: %s", fp.name, exc)
    except Exception as exc:
        app.logger.warning("Worship seed failed: %s", exc)


def _invalidate_worship_cache(scope: str | None = None) -> None:
    """Drop the cached song list for one scope (default: the current one)."""
    global _worship_songs_cache
    if scope is None:
        try:
            scope = _current_worship_scope()
        except Exception:
            scope = None
    if scope is None:
        _worship_songs_cache = {}
    else:
        _worship_songs_cache.pop(scope, None)


def _worship_song_sort_key(song: dict) -> tuple[str, str, str]:
    normalized = normalize_worship_song(song)
    return (
        normalized.get("title", "").lower(),
        normalized.get("artist", "").lower(),
        normalized.get("version", "").lower(),
        normalized.get("id", "").lower(),
    )


def list_worship_songs() -> list[dict]:
    """List all worship songs sorted by title. Cached for _WORSHIP_CACHE_TTL seconds."""
    import time
    now = time.monotonic()
    scope = _current_worship_scope()
    entry = _worship_songs_cache.get(scope)
    if entry and (now - entry["ts"]) < _WORSHIP_CACHE_TTL:
        return entry["data"]

    if db:
        try:
            by_id = {}
            for ref in _worship_song_refs_for_read():
                for doc in ref.order_by("title").stream():
                    data = doc.to_dict()
                    if data and data.get("id"):
                        by_id[data["id"]] = data
            results = list(by_id.values())
            results.sort(key=_worship_song_sort_key)
            _worship_songs_cache[scope] = {"data": results, "ts": now}
            return results
        except Exception as exc:
            app.logger.warning("list_worship_songs Firestore error: %s", exc)
            if not _is_local_storage_allowed():
                raise RuntimeError("Firestore query failed for list_worship_songs in production.") from exc
    if not _is_local_storage_allowed():
        raise RuntimeError("Firestore not available to list songs in production. Local fallback is disabled.")
    songs_folder = Path(app.root_path) / "songs"
    songs = []
    if songs_folder.is_dir():
        for fp in sorted(songs_folder.glob("*.json")):
            try:
                with open(fp, "r", encoding="utf-8") as f:
                    songs.append(json.load(f))
            except Exception:
                pass
    songs.sort(key=_worship_song_sort_key)
    _worship_songs_cache[scope] = {"data": songs, "ts": now}
    return songs


def get_worship_song(song_id: str) -> dict | None:
    """Load one song by id. Firestore first, file fallback."""
    firestore_client, _ = init_firebase()
    if firestore_client is not None:
        try:
            for ref in _worship_song_refs_for_read(firestore_client):
                doc = ref.document(song_id).get()
                if doc.exists:
                    return doc.to_dict()
        except Exception as exc:
            app.logger.warning("get_worship_song(%s) Firestore error: %s", song_id, exc)
            if not _is_local_storage_allowed():
                raise RuntimeError(f"Firestore read failed for song {song_id} in production.") from exc
        else:
            # A successful query with no matching document means the song is new;
            # it does not mean Firestore is unavailable.
            if not _is_local_storage_allowed():
                return None
    if not _is_local_storage_allowed():
        app.logger.error(
            "get_worship_song(%s) has no Firestore client in pid=%s: %s",
            song_id, os.getpid(), firebase_init_diagnostic(),
        )
        raise RuntimeError(f"Firestore not available to get song {song_id} in production. Local fallback is disabled.")
    fp = Path(app.root_path) / "songs" / f"{song_id}.json"
    if fp.exists():
        try:
            with open(fp, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    return None


def _canonical_part_key(name: str | None) -> str:
    raw = str(name or "").strip().lower()
    if not raw:
        return ""
    raw = raw.replace("&", "and")
    raw = re.sub(r"[^a-z0-9]+", "_", raw).strip("_")
    alias_map = {
        "prechorus": "pre_chorus",
        "pre_chorus_1": "pre_chorus1",
        "pre_chorus_2": "pre_chorus2",
        "refrain": "chorus",
        "verse_1": "verse1",
        "verse_2": "verse2",
        "verse_3": "verse3",
        "verse_4": "verse4",
        "chorus_1": "chorus",
        "chorus_2": "chorus2",
        "chorus_3": "chorus3",
        "bridge_1": "bridge1",
        "bridge_2": "bridge2",
        "tag_1": "tag1",
        "tag_2": "tag2",
        "ending": "outro",
    }
    if raw in alias_map:
        return alias_map[raw]
    numbered_match = re.match(r"^(verse|chorus|bridge|tag)_(\d+)$", raw)
    if numbered_match:
        return f"{numbered_match.group(1)}{numbered_match.group(2)}"
    return raw


_REUSABLE_LYRIC_SHEET_PART_PREFIXES = ("chorus", "pre_chorus", "bridge", "tag")


def _is_reusable_part(part_name: str) -> bool:
    return any(part_name == p or part_name.startswith(p) for p in _REUSABLE_LYRIC_SHEET_PART_PREFIXES)


def _worship_part_label(part_name: str) -> str:
    label = str(part_name or "").replace("_", " ")
    label = re.sub(r"([a-zA-Z])(\d+)", r"\1 \2", label)
    return label.title()


def normalize_worship_song(song: dict) -> dict:
    normalized = dict(song or {})
    title = str(normalized.get("title") or "").strip()
    normalized["title"] = title
    normalized["artist"] = str(normalized.get("artist") or "").strip()
    normalized["version"] = str(normalized.get("version") or "").strip()
    normalized["key"] = str(normalized.get("key") or "").strip()
    normalized["type"] = str(normalized.get("type") or "song").strip() or "song"
    normalized["background"] = str(normalized.get("background") or "").strip()
    song_id = str(normalized.get("id") or "").strip()
    if not song_id:
        song_id = re.sub(r"[^a-z0-9]+", "-", title.lower()).strip("-") or "untitled-song"
    normalized["id"] = song_id

    raw_parts = normalized.get("parts") if isinstance(normalized.get("parts"), dict) else {}
    parts: dict[str, list[str]] = {}
    for raw_name, raw_lines in raw_parts.items():
        part_name = _canonical_part_key(raw_name)
        if not part_name:
            continue
        if isinstance(raw_lines, list):
            cleaned_lines = [str(line).strip() for line in raw_lines if str(line).strip()]
        else:
            cleaned_lines = [str(raw_lines).strip()] if str(raw_lines).strip() else []
        if not cleaned_lines:
            continue
        if part_name in parts:
            existing = parts[part_name]
            if existing != cleaned_lines:
                merged = existing + [line for line in cleaned_lines if line not in existing]
                parts[part_name] = merged
        else:
            parts[part_name] = cleaned_lines

    raw_arrangement = normalized.get("arrangement") if isinstance(normalized.get("arrangement"), list) else []
    arrangement: list[str] = []
    for raw_part in raw_arrangement:
        part_name = _canonical_part_key(raw_part)
        if part_name and part_name in parts:
            arrangement.append(part_name)

    if not arrangement:
        arrangement = list(parts.keys())

    normalized["parts"] = parts
    normalized["arrangement"] = arrangement
    return normalized


def build_lyric_sheet_blocks(song: dict) -> list[dict]:
    normalized = normalize_worship_song(song)
    parts = normalized.get("parts", {}) or {}
    arrangement = normalized.get("arrangement", []) or []
    blocks: list[dict] = []
    seen_parts: set[str] = set()
    for part_name in arrangement:
        lines = parts.get(part_name, [])
        if not lines:
            continue
        repeated = part_name in seen_parts
        reference_only = repeated and _is_reusable_part(part_name)
        blocks.append(
            {
                "part": part_name,
                "label": _worship_part_label(part_name),
                "lines": [] if reference_only else list(lines),
                "reference_only": reference_only,
            }
        )
        seen_parts.add(part_name)
    return blocks


def save_worship_song(song: dict) -> None:
    """Persist a song. Firestore when available, local file fallback."""
    song = normalize_worship_song(song)
    song_id = song["id"]
    firestore_client, _ = init_firebase()
    if firestore_client is not None:
        try:
            song["worship_scope"] = _current_worship_scope()
            _worship_songs_ref(firestore_client).document(song_id).set(song)
            _invalidate_worship_cache()
            return
        except Exception as exc:
            app.logger.warning("save_worship_song(%s) Firestore error: %s", song_id, exc)
            if not _is_local_storage_allowed():
                raise RuntimeError(f"Firestore write failed for song {song_id} in production.") from exc
    if not _is_local_storage_allowed():
        app.logger.error(
            "save_worship_song(%s) has no Firestore client in pid=%s: %s",
            song_id, os.getpid(), firebase_init_diagnostic(),
        )
        raise RuntimeError(f"Firestore not available to save song {song_id} in production. Local fallback is disabled.")
    songs_folder = Path(app.root_path) / "songs"
    songs_folder.mkdir(exist_ok=True)
    with open(songs_folder / f"{song_id}.json", "w", encoding="utf-8") as f:
        json.dump(song, f, indent=2, ensure_ascii=False)
    _invalidate_worship_cache()


_WORSHIP_PENDING_COLLECTION = "worship_imports_pending"
_WORSHIP_PENDING_TTL = timedelta(minutes=30)
_local_pending_worship_imports: dict[str, dict] = {}
_local_pending_worship_lock = threading.Lock()


def _store_pending_worship_song(song: dict, used_fallback: bool, fallback_reason: str) -> str:
    token = secrets.token_urlsafe(24)
    payload = {
        "song": normalize_worship_song(song),
        "used_fallback": bool(used_fallback),
        "fallback_reason": str(fallback_reason or ""),
        "owner": str(session.get("user_email") or ""),
        "scope": _current_worship_scope(),
        "expires_at": datetime.now(timezone.utc) + _WORSHIP_PENDING_TTL,
    }
    if db:
        try:
            firestore_payload = dict(payload)
            firestore_payload["expireAt"] = firestore_payload.pop("expires_at")
            db.collection(_WORSHIP_PENDING_COLLECTION).document(token).set(firestore_payload)
            return token
        except Exception as exc:
            app.logger.warning("Pending worship import write failed: %s", exc)
            if not _is_local_storage_allowed():
                raise RuntimeError("Could not preserve the pending worship import.") from exc
    if not _is_local_storage_allowed():
        raise RuntimeError("Could not preserve the pending worship import.")
    with _local_pending_worship_lock:
        _local_pending_worship_imports[token] = payload
    return token


def _load_pending_worship_song(token: str) -> dict | None:
    token = str(token or "").strip()
    if not token:
        return None
    payload = None
    if db:
        try:
            doc = db.collection(_WORSHIP_PENDING_COLLECTION).document(token).get()
            payload = doc.to_dict() if doc.exists else None
        except Exception as exc:
            app.logger.warning("Pending worship import read failed: %s", exc)
            if not _is_local_storage_allowed():
                raise RuntimeError("Could not load the pending worship import.") from exc
    elif not _is_local_storage_allowed():
        raise RuntimeError("Could not load the pending worship import.")
    if payload is None and _is_local_storage_allowed():
        with _local_pending_worship_lock:
            payload = _local_pending_worship_imports.get(token)
    if not isinstance(payload, dict):
        return None
    if payload.get("owner") != str(session.get("user_email") or ""):
        return None
    if payload.get("scope") != _current_worship_scope():
        return None
    expires_at = payload.get("expireAt") or payload.get("expires_at")
    if isinstance(expires_at, datetime):
        if expires_at.tzinfo is None:
            expires_at = expires_at.replace(tzinfo=timezone.utc)
    if isinstance(expires_at, datetime) and expires_at < datetime.now(timezone.utc):
        _delete_pending_worship_song(token)
        return None
    return payload


def _delete_pending_worship_song(token: str) -> None:
    token = str(token or "").strip()
    if not token:
        return
    if db:
        try:
            db.collection(_WORSHIP_PENDING_COLLECTION).document(token).delete()
        except Exception as exc:
            app.logger.warning("Pending worship import cleanup failed: %s", exc)
    with _local_pending_worship_lock:
        _local_pending_worship_imports.pop(token, None)


def delete_worship_song(song_id: str) -> bool:
    """Delete a song from Firestore and local file if present. Returns True if anything was deleted."""
    deleted = False
    if db:
        try:
            collection_refs = [_worship_songs_ref()]
            if _current_worship_scope() == _DEFAULT_WORSHIP_SCOPE:
                collection_refs.append(_legacy_worship_songs_ref())
            for collection_ref in collection_refs:
                ref = collection_ref.document(song_id)
                if ref.get().exists:
                    ref.delete()
                    deleted = True
                    continue
                for doc in collection_ref.stream():
                    data = doc.to_dict() or {}
                    if str(data.get("id") or "").strip() == song_id:
                        doc.reference.delete()
                        deleted = True
                        break
        except Exception as exc:
            app.logger.warning("delete_worship_song(%s) Firestore error: %s", song_id, exc)
            if not _is_local_storage_allowed():
                raise RuntimeError(f"Firestore delete failed for song {song_id} in production.") from exc
    if not _is_local_storage_allowed():
        if not db:
            raise RuntimeError(f"Firestore not available to delete song {song_id} in production. Local fallback is disabled.")
    if _is_local_storage_allowed():
        fp = Path(app.root_path) / "songs" / f"{song_id}.json"
        if fp.exists():
            try:
                fp.unlink()
                deleted = True
            except Exception:
                pass
    if deleted:
        _invalidate_worship_cache()
    return deleted


def _delete_all_worship_library_data() -> dict:
    """Delete songs and saved setlists for the active worship scope."""
    deleted = {"songs": 0, "setlists": 0}
    if db:
        try:
            for collection_ref in _worship_song_refs_for_read():
                for doc in collection_ref.stream():
                    doc.reference.delete()
                    deleted["songs"] += 1
            for collection_ref in _worship_setlist_refs_for_read():
                for doc in collection_ref.stream():
                    doc.reference.delete()
                    deleted["setlists"] += 1
        except Exception as exc:
            app.logger.warning("_delete_all_worship_library_data Firestore error: %s", exc)
    if not db:
        for folder_name, key in (("songs", "songs"), ("setlists", "setlists")):
            folder = Path(app.root_path) / folder_name
            if not folder.is_dir():
                continue
            for fp in folder.glob("*.json"):
                try:
                    fp.unlink()
                    deleted[key] += 1
                except Exception as exc:
                    app.logger.warning("Could not delete %s: %s", fp, exc)
    if deleted["songs"] or deleted["setlists"]:
        _invalidate_worship_cache()
    return deleted


def _remove_song_from_worship_setlists(song_id: str) -> int:
    """Remove a deleted song id from saved setlists. Empty setlists are deleted."""
    changed = 0
    if db:
        try:
            for collection_ref in _worship_setlist_refs_for_read():
                for doc in collection_ref.stream():
                    data = doc.to_dict() or {}
                    songs = data.get("songs")
                    if not isinstance(songs, list) or song_id not in songs:
                        continue
                    notes = data.get("notes") if isinstance(data.get("notes"), dict) else {}
                    data["songs"] = [sid for sid in songs if sid != song_id]
                    data["notes"] = {sid: note for sid, note in notes.items() if sid != song_id}
                    if data["songs"]:
                        doc.reference.set(data)
                    else:
                        doc.reference.delete()
                    changed += 1
        except Exception as exc:
            app.logger.warning("_remove_song_from_worship_setlists(%s) Firestore error: %s", song_id, exc)

    setlists_dir = Path(app.root_path) / "setlists"
    if setlists_dir.is_dir():
        for fp in setlists_dir.glob("*.json"):
            try:
                with open(fp, "r", encoding="utf-8") as f:
                    data = json.load(f)
                songs = data.get("songs")
                if not isinstance(songs, list) or song_id not in songs:
                    continue
                notes = data.get("notes") if isinstance(data.get("notes"), dict) else {}
                data["songs"] = [sid for sid in songs if sid != song_id]
                data["notes"] = {sid: note for sid, note in notes.items() if sid != song_id}
                if data["songs"]:
                    with open(fp, "w", encoding="utf-8") as f:
                        json.dump(data, f, indent=2)
                else:
                    fp.unlink()
                changed += 1
            except Exception as exc:
                app.logger.warning("_remove_song_from_worship_setlists(%s) file error: %s", song_id, exc)
    return changed


def _worship_wants_json_response() -> bool:
    return request.accept_mimetypes.best == "application/json" or request.headers.get("X-Requested-With") == "fetch"


def _worship_setlist_id(date_label: str, name: str = "") -> str:
    name_slug = _slugify_worship_token(name)
    return f"{date_label}-{name_slug}" if name_slug else date_label


def _worship_timestamp() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat()


def _normalize_worship_setlist(data: dict, fallback_id: str = "") -> dict:
    date_label = str(data.get("date") or fallback_id[:10] or "").strip()
    songs = data.get("songs") if isinstance(data.get("songs"), list) else []
    notes = data.get("notes") if isinstance(data.get("notes"), dict) else {}
    name = str(data.get("name") or "").strip()
    setlist_id = str(data.get("id") or fallback_id or _worship_setlist_id(date_label, name)).strip()
    return {
        "id": setlist_id,
        "date": date_label,
        "name": name,
        "songs": songs,
        "notes": notes,
        "created_at": data.get("created_at") or data.get("createdAt") or "",
        "updated_at": data.get("updated_at") or data.get("updatedAt") or "",
        "created_by": data.get("created_by") or data.get("createdBy") or "",
        "updated_by": data.get("updated_by") or data.get("updatedBy") or "",
        "song_count": int(data.get("song_count") or len(songs)),
        "worship_scope": data.get("worship_scope") or data.get("worshipScope") or _current_worship_scope(),
    }


def _worship_setlist_visible_in_scope(data: dict) -> bool:
    scope = data.get("worship_scope") or data.get("worshipScope") or _DEFAULT_WORSHIP_SCOPE
    current = _current_worship_scope()
    return scope == current or (current == _DEFAULT_WORSHIP_SCOPE and scope in ("", _DEFAULT_WORSHIP_SCOPE))


def _valid_worship_setlist_id(setlist_id: str) -> bool:
    return bool(setlist_id) and ".." not in setlist_id and "/" not in setlist_id and "\\" not in setlist_id


def _get_worship_setlist(setlist_id: str) -> dict | None:
    if not _valid_worship_setlist_id(setlist_id):
        return None
    if db:
        try:
            for collection_ref in _worship_setlist_refs_for_read():
                doc = collection_ref.document(setlist_id).get()
                if doc.exists:
                    return _normalize_worship_setlist(doc.to_dict() or {}, doc.id)
        except Exception as exc:
            app.logger.warning("_get_worship_setlist(%s) Firestore error: %s", setlist_id, exc)
            if not _is_local_storage_allowed():
                raise RuntimeError(f"Firestore read failed for setlist {setlist_id} in production.") from exc
        else:
            if not _is_local_storage_allowed():
                return None
    if not _is_local_storage_allowed():
        raise RuntimeError(f"Firestore not available to get setlist {setlist_id} in production. Local fallback is disabled.")
    fp = Path(app.root_path) / "setlists" / f"{setlist_id}.json"
    if fp.exists():
        try:
            with open(fp, "r", encoding="utf-8") as f:
                data = json.load(f)
            if _worship_setlist_visible_in_scope(data):
                return _normalize_worship_setlist(data, fp.stem)
        except Exception:
            pass
    return None


def _make_unique_worship_setlist_id(date_label: str, name: str, existing_id: str = "") -> str:
    base = _worship_setlist_id(date_label, name)
    candidate = base
    suffix = 2
    while candidate != existing_id and _get_worship_setlist(candidate):
        candidate = f"{base}-{suffix}"
        suffix += 1
    return candidate


def _persist_worship_setlist(data: dict, previous_id: str = "") -> bool:
    setlist_id = data["id"]
    existed = bool(_get_worship_setlist(setlist_id))
    if db:
        try:
            scoped_ref = _worship_setlists_ref().document(setlist_id)
            scoped_ref.set(data)
            if previous_id and previous_id != setlist_id:
                _worship_setlists_ref().document(previous_id).delete()
                if _current_worship_scope() == _DEFAULT_WORSHIP_SCOPE:
                    _legacy_worship_setlists_ref().document(previous_id).delete()
        except Exception as exc:
            app.logger.warning("_persist_worship_setlist Firestore error: %s", exc)
            if not _is_local_storage_allowed():
                raise RuntimeError(f"Firestore save failed for setlist {setlist_id} in production.") from exc
    else:
        if not _is_local_storage_allowed():
            raise RuntimeError(f"Firestore not available to save setlist {setlist_id} in production. Local fallback is disabled.")

    if _is_local_storage_allowed():
        try:
            setlists_dir = Path(app.root_path) / "setlists"
            setlists_dir.mkdir(exist_ok=True)
            with open(setlists_dir / f"{setlist_id}.json", "w", encoding="utf-8") as f:
                json.dump(data, f, indent=2)
            if previous_id and previous_id != setlist_id:
                old_fp = setlists_dir / f"{previous_id}.json"
                if old_fp.exists():
                    old_fp.unlink()
        except Exception as exc:
            app.logger.warning("_persist_worship_setlist file error: %s", exc)
    return existed


def _delete_worship_setlist(setlist_id: str) -> bool:
    if not _valid_worship_setlist_id(setlist_id):
        return False
    deleted = False
    if db:
        try:
            ref = _worship_setlists_ref().document(setlist_id)
            if ref.get().exists:
                deleted = True
            ref.delete()
            if _current_worship_scope() == _DEFAULT_WORSHIP_SCOPE:
                legacy_ref = _legacy_worship_setlists_ref().document(setlist_id)
                if legacy_ref.get().exists:
                    deleted = True
                legacy_ref.delete()
        except Exception as exc:
            app.logger.warning("_delete_worship_setlist Firestore error: %s", exc)
            if not _is_local_storage_allowed():
                raise RuntimeError(f"Firestore delete failed for setlist {setlist_id} in production.") from exc
    else:
        if not _is_local_storage_allowed():
            raise RuntimeError(f"Firestore not available to delete setlist {setlist_id} in production. Local fallback is disabled.")

    if _is_local_storage_allowed():
        fp = Path(app.root_path) / "setlists" / f"{setlist_id}.json"
        if fp.exists():
            try:
                fp.unlink()
                deleted = True
            except Exception:
                pass
    return deleted


def _touch_worship_songs_last_used(song_ids: list[str], date_label: str) -> None:
    for song_id in dict.fromkeys(song_ids):
        song = get_worship_song(song_id)
        if not song:
            continue
        song["last_used"] = date_label
        save_worship_song(song)


def _load_recent_setlists() -> list[dict]:
    """Return all setlists, newest first."""
    if db:
        try:
            by_id = {}
            for collection_ref in _worship_setlist_refs_for_read():
                for doc in collection_ref.stream():
                    data = doc.to_dict()
                    if data and data.get("date") and isinstance(data.get("songs"), list):
                        normalized = _normalize_worship_setlist(data, doc.id)
                        by_id[normalized["id"]] = normalized
            results = list(by_id.values())
            results.sort(key=lambda x: (x["date"], x["name"], x["id"]), reverse=True)
            return results
        except Exception as exc:
            app.logger.warning("_load_recent_setlists Firestore error: %s", exc)
    setlists_dir = Path(app.root_path) / "setlists"
    if not setlists_dir.is_dir():
        return []
    results = []
    for fp in sorted(setlists_dir.glob("*.json"), reverse=True):
        try:
            with open(fp, "r", encoding="utf-8") as f:
                data = json.load(f)
            if data.get("date") and isinstance(data.get("songs"), list) and _worship_setlist_visible_in_scope(data):
                results.append(_normalize_worship_setlist(data, fp.stem))
        except Exception:
            pass
    results.sort(key=lambda x: (x["date"], x["name"], x["id"]), reverse=True)
    return results


_CONTINUATION_STARTS = (
    "and ",
    "but ",
    "or ",
    "so ",
    "for ",
    "nor ",
    "yet ",
    "because ",
    "cause ",
    "'cause ",
    "inside ",
    "with ",
    "to ",
    "of ",
    "in ",
    "on ",
)


def _line_continues_to_next(line: str) -> bool:
    text = (line or "").strip().lower()
    if not text:
        return False
    if text.endswith((",", ";", ":", "-", "(", "...", "—")):
        return True
    if text.endswith(("and", "or", "but", "so", "cause", "because", "to", "of", "in", "with")):
        return True
    return False


def _line_starts_as_continuation(line: str) -> bool:
    text = (line or "").strip().lower()
    return any(text.startswith(prefix) for prefix in _CONTINUATION_STARTS)


def _is_protected_phrase_pair(prev_line: str, next_line: str) -> bool:
    prev = (prev_line or "").strip().lower()
    nxt = (next_line or "").strip().lower()
    protected_pairs = [
        ("cause you've got a lion", "inside of those lungs"),
        ("get up and", "praise the lord"),
    ]
    return any(prev.endswith(a) and nxt.startswith(b) for a, b in protected_pairs)


def _is_phrase_boundary(line: str) -> bool:
    text = (line or "").strip()
    if not text:
        return False
    return text.endswith((".", "!", "?", ":", ";"))


def _join_required(prev_line: str, next_line: str) -> bool:
    return (
        _is_protected_phrase_pair(prev_line, next_line)
        or _line_continues_to_next(prev_line)
        or _line_starts_as_continuation(next_line)
    )




def _split_at_midpoint(lines: list[str]) -> tuple[list[str], list[str]]:
    """Split near midpoint, preferring phrase boundaries."""
    n = len(lines)
    mid = n // 2
    for offset in range(n):
        for sign in (0, 1, -1):
            idx = mid + sign * offset
            if idx <= 0 or idx >= n:
                continue
            if not _join_required(lines[idx - 1], lines[idx]):
                return lines[:idx], lines[idx:]
    return lines[:mid], lines[mid:]  # fallback: hard split


def chunk_lines(lines: list[str]) -> list[dict]:
    """
    Lyric-aware chunking for PPTX slides.

    Rules:
      1. Blank lines are hard slide-break hints.
      2. Groups of ≤4 lines → one slide, font_size=48.
      3. A group of 5 lines → one slide, font_size=42.
      4. Groups of 6+ lines → split near midpoint (phrase-boundary-aware),
         recursively until ≤5 lines; font_size derived from final chunk size.
      5. Never splits mid-phrase (_join_required respected throughout).
    """
    groups: list[list[str]] = []
    current: list[str] = []
    for raw in lines:
        line = raw.strip() if isinstance(raw, str) else ""
        if not line:
            if current:
                groups.append(current)
                current = []
        else:
            current.append(line)
    if current:
        groups.append(current)

    def _chunks(grp: list[str]) -> list[list[str]]:
        if len(grp) <= 5:
            return [grp]
        left, right = _split_at_midpoint(grp)
        return _chunks(left) + _chunks(right)

    slides: list[dict] = []
    for group in groups:
        for chunk in _chunks(group):
            n = len(chunk)
            font_size = 48 if n <= 4 else 42 if n == 5 else 38
            slides.append({"lines": chunk, "font_size": font_size})
    return slides


_TYPE_BG = {
    "song":      RGBColor(28, 28, 28),
    "scripture": RGBColor(10, 20, 55),
    "prayer":    RGBColor(10, 40, 20),
    "reading":   RGBColor(10, 40, 20),
}
_DEFAULT_BG = RGBColor(28, 28, 28)

_BG_DIR = pathlib.Path(__file__).parent / "static" / "worship" / "backgrounds"
_DEFAULT_BACKGROUNDS = {
    "song":      "deep-blue-abstract.png",
    "scripture": "parchment-texture.png",
    "prayer":    "dark-green-abstract.png",
    "reading":   "soft-clouds.png",
}
_FALLBACK_BG = "dark-gradient.png"
_bg_config_cache: dict | None = None


def _load_bg_config() -> dict:
    global _bg_config_cache
    if _bg_config_cache is not None:
        return _bg_config_cache
    try:
        with open(_BG_DIR / "backgrounds.json", "r", encoding="utf-8") as f:
            _bg_config_cache = json.load(f)
    except Exception:
        _bg_config_cache = {}
    return _bg_config_cache


def _resolve_bg(item_type: str, song_bg: str | None) -> tuple:
    """Return (Path | None, config dict) for the best available background."""
    cfg = _load_bg_config()
    for filename in (song_bg, _DEFAULT_BACKGROUNDS.get(item_type), _FALLBACK_BG):
        if not filename:
            continue
        path = _BG_DIR / filename
        if path.exists():
            defaults = {"font_color": [255, 255, 255], "overlay": False, "overlay_opacity": 0}
            return path, {**defaults, **cfg.get(filename, {})}
    return None, {"font_color": [255, 255, 255], "overlay": False, "overlay_opacity": 0}


def _apply_image_background(slide, img_path) -> None:
    """Add image as full-slide background, behind all other shapes."""
    pic = slide.shapes.add_picture(str(img_path), Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    sp_tree = slide.shapes._spTree
    sp_tree.remove(pic._element)
    sp_tree.insert(2, pic._element)


def _add_overlay_rect(slide, left: float, top: float, width: float, height: float, opacity: float) -> None:
    """Add a semi-transparent black rectangle matching textbox bounds."""
    from lxml import etree
    from pptx.oxml.ns import qn
    rect = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0, 0, 0)
    rect.line.fill.background()
    solidFill = rect.fill._xPr.find('.//' + qn('a:solidFill'))
    if solidFill is not None:
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is not None:
            etree.SubElement(srgbClr, qn('a:alpha')).set('val', str(int(opacity * 100000)))


def apply_dark_background(slide, rgb: RGBColor = _DEFAULT_BG) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb


def add_centered_textbox(slide, text: str, top: float, height: float, font_size: int, bold: bool,
                          font_color: RGBColor = None, overlay: bool = False, overlay_opacity: float = 0.5,
                          left: float = 0.5, width: float = 12.333):
    if font_color is None:
        font_color = RGBColor(255, 255, 255)
    if overlay:
        _add_overlay_rect(slide, left, top, width, height, overlay_opacity)
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = "Arial"
    run.font.color.rgb = font_color
    return box


def add_link_footer(slide, label: str = "faithsparksprintables.com", url: str = "https://faithsparksprintables.com"):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(7.02), Inches(12.5), Inches(0.18))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(9)
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(235, 235, 235)
    run.hyperlink.address = url
    return box


def create_divider_slide(prs: Presentation, title: str, artist: str, key: str, item_type: str = "song", song_bg: str = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    img_path, cfg = _resolve_bg(item_type, song_bg)
    if img_path:
        _apply_image_background(slide, img_path)
    else:
        apply_dark_background(slide, _TYPE_BG.get(item_type, _DEFAULT_BG))
    font_color = RGBColor(*cfg["font_color"])
    overlay = cfg.get("overlay", False)
    overlay_opacity = cfg.get("overlay_opacity", 0.5)
    add_centered_textbox(slide, title, top=2.2, height=1.65, font_size=60, bold=True,
                          font_color=font_color, overlay=overlay, overlay_opacity=overlay_opacity,
                          left=1.5, width=10.333)
    subtitle_parts = []
    if artist:
        subtitle_parts.append(artist)
    if key:
        subtitle_parts.append(f"Key: {key}")
    if subtitle_parts:
        add_centered_textbox(slide, " | ".join(subtitle_parts), top=4.05, height=0.85, font_size=26, bold=False,
                              font_color=font_color, overlay=overlay, overlay_opacity=overlay_opacity,
                              left=1.5, width=10.333)
    add_link_footer(slide)
    return slide


def _add_part_label(slide, label: str, font_color: RGBColor) -> None:
    if not label:
        return
    box = slide.shapes.add_textbox(Inches(0.42), Inches(0.34), Inches(4.0), Inches(0.3))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = label.upper()
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.name = "Arial"
    r = int(font_color[0] * 0.4)
    g = int(font_color[1] * 0.4)
    b = int(font_color[2] * 0.4)
    run.font.color.rgb = RGBColor(r, g, b)


def create_content_slide(prs: Presentation, lines: list[str], item_type: str = "song", song_bg: str = None, font_size: int = 48, part_label: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    img_path, cfg = _resolve_bg(item_type, song_bg)
    if img_path:
        _apply_image_background(slide, img_path)
    else:
        apply_dark_background(slide, _TYPE_BG.get(item_type, _DEFAULT_BG))
    font_color = RGBColor(*cfg["font_color"])
    overlay = cfg.get("overlay", False)
    overlay_opacity = cfg.get("overlay_opacity", 0.5)
    _add_part_label(slide, part_label, font_color)
    add_centered_textbox(slide, "\n".join(lines), top=1.25, height=5.0, font_size=font_size, bold=True,
                          font_color=font_color, overlay=overlay, overlay_opacity=overlay_opacity)
    add_link_footer(slide)
    return slide



os.makedirs("output", exist_ok=True)
os.makedirs("output/thumbs", exist_ok=True)
os.makedirs("output/packs", exist_ok=True)

# --- Theming ---
# theme helpers moved to yourapp.services.themes

@app.route("/login/google/start")
def start_google_login():
    nxt = request.args.get("next")
    # Default to /browse if missing/unsafe
    if not is_safe_url(nxt or ""):
        nxt = url_for("browse")
    session["post_login_next"] = nxt
    return redirect(url_for("google.login", next=nxt))

@app.route("/oauth/finish")
def oauth_finish():
    # by now google.authorized is True; @before_request already hydrated user_info
    nxt = session.pop("post_login_next", None) or request.args.get("next")

    email = session.get("user_email")
    if email:
        try:
            analytics_svc.record_login(email)
        except Exception:
            app.logger.debug("Failed to record login", exc_info=True)

    if nxt and is_safe_url(nxt):
        return redirect(nxt)

    # Safe fallback if public blueprint isn't available for any reason
    if "public.index" in app.view_functions:
        return redirect(url_for("public.index"))
    return redirect(url_for("browse"))


@app.route('/admin/theme/preview', methods=['POST'])
@admin_required
def admin_theme_preview():
    from faithsparks.views.admin.theme import admin_theme_preview as _impl
    return _impl()

    

"""Collections helpers moved to yourapp.services.collections"""

 

# storage helpers moved to yourapp.services.storage

# --- Routes ---

# Mount public routes blueprint
try:
    from faithsparks.views.public import bp as public_bp
    app.register_blueprint(public_bp)
except Exception:
    pass

# Family Bible Bee room-code game
try:
    from faithsparks.views.bible_bee import bp as bible_bee_bp
    app.register_blueprint(bible_bee_bp)
except Exception:
    app.logger.exception("Family Bible Bee routes could not be registered")

# Group game night room-code games
try:
    from faithsparks.views.act_it_out import bp as act_it_out_bp
    app.register_blueprint(act_it_out_bp)
except Exception:
    app.logger.exception("Act It Out routes could not be registered")

@app.route("/logout")
def logout():
    session.clear()
    if "public.index" in app.view_functions:
        return redirect(url_for("public.index"))
    return redirect(url_for("browse"))

from pathlib import Path
from tempfile import NamedTemporaryFile
import json


@app.route('/worship', methods=['GET'])
@login_required
def worship():
    _seed_worship_from_files()
    songs = list_worship_songs()
    setlists = _load_recent_setlists()
    church_context = _current_worship_church_context()
    return render_template('worship.html', songs=songs, setlists=setlists, worship_church=church_context["current"], worship_churches=church_context["churches"])


@app.route("/worship/cleanup-duplicates", methods=["POST"])
@login_required
def worship_cleanup_duplicates():
    """Delete duplicate songs (same title+artist), keeping the one with the most lyric content."""
    songs = list_worship_songs()

    # Group by (title, artist, version) — different versions are intentional arrangements,
    # never merge them. Only deduplicate true re-imports of the same arrangement.
    groups: dict[tuple, list[dict]] = {}
    for song in songs:
        normalized = normalize_worship_song(song)
        key = (
            normalized.get("title", "").lower().strip(),
            normalized.get("artist", "").lower().strip(),
            normalized.get("version", "").lower().strip(),
        )
        groups.setdefault(key, []).append(normalized)

    deleted_ids: list[str] = []
    kept_ids: list[str] = []

    def _parse_quality_score(s: dict) -> int:
        """Higher = better structured parse. Penalises bridge2/bridge3/verse9-style fragmentation."""
        parts = s.get("parts") or {}
        total_lines = sum(len(lines) for lines in parts.values() if isinstance(lines, list))
        # Count fragmented numbered parts (bridge2, verse9, etc.) — each costs 20 points
        frag_penalty = sum(
            20 for name in parts
            if re.match(r'^(bridge|verse|chorus|tag|pre_chorus)\d{1,2}$', name)
            and not name in ("verse1", "verse2", "verse3", "chorus", "chorus2", "bridge", "pre_chorus", "tag", "intro", "outro")
        )
        return total_lines - frag_penalty

    for key, group in groups.items():
        if len(group) <= 1:
            continue

        group.sort(key=_parse_quality_score, reverse=True)
        best = group[0]
        kept_ids.append(best["id"])

        for dupe in group[1:]:
            if dupe["id"] != best["id"]:
                delete_worship_song(dupe["id"])
                deleted_ids.append(dupe["id"])

    if deleted_ids:
        flash(
            f"Removed {len(deleted_ids)} duplicate{'s' if len(deleted_ids) != 1 else ''}: {', '.join(deleted_ids)}.",
            "success",
        )
    else:
        flash("No duplicates found.", "info")

    return redirect(url_for("worship"))


@app.route("/worship/church/create", methods=["POST"])
@login_required
def worship_church_create():
    if not db:
        flash("Church sharing needs Firestore.", "warning")
        return redirect(url_for("worship"))
    name = request.form.get("church_name", "").strip()
    if not name:
        flash("Church name is required.", "warning")
        return redirect(url_for("worship"))
    church_id = _slugify_worship_token(name) or f"church-{secrets.token_hex(3)}"
    base_id = church_id
    suffix = 2
    try:
        while db.collection(_WORSHIP_CHURCH_COLLECTION).document(church_id).get().exists:
            church_id = f"{base_id}-{suffix}"
            suffix += 1
        invite_code = _make_worship_invite_code(name)
        church_ref = db.collection(_WORSHIP_CHURCH_COLLECTION).document(church_id)
        church_ref.set(
            {
                "id": church_id,
                "name": name,
                "invite_code": invite_code,
                "created_by": session.get("user_email"),
                "created_at": firestore.SERVER_TIMESTAMP,
                "updated_at": firestore.SERVER_TIMESTAMP,
            }
        )
        church_ref.collection("members").document(session["user_email"]).set(
            {"email": session["user_email"], "role": "owner", "joined_at": firestore.SERVER_TIMESTAMP},
            merge=True,
        )
        _set_current_worship_church(church_id)
        _invalidate_worship_cache()
        flash(f"Created {name}. Invite code: {invite_code}", "success")
    except Exception as exc:
        app.logger.warning("worship_church_create error: %s", exc)
        flash("We couldn't create that church. Please try again.", "error")
    return redirect(url_for("worship"))


@app.route("/worship/church/join", methods=["POST"])
@login_required
def worship_church_join():
    if not db:
        flash("Church sharing needs Firestore.", "warning")
        return redirect(url_for("worship"))
    invite_code = _normalize_worship_invite_code(request.form.get("invite_code", ""))
    if not invite_code:
        flash("Enter a church invite code.", "warning")
        return redirect(url_for("worship"))
    try:
        matches = (
            db.collection(_WORSHIP_CHURCH_COLLECTION)
            .where(filter=firestore.FieldFilter("invite_code", "==", invite_code))
            .limit(1)
            .stream()
        )
        church_doc = next(matches, None)
        if not church_doc:
            flash("That church code was not found.", "error")
            return redirect(url_for("worship"))
        church = church_doc.to_dict() or {}
        church_doc.reference.collection("members").document(session["user_email"]).set(
            {"email": session["user_email"], "role": "member", "joined_at": firestore.SERVER_TIMESTAMP},
            merge=True,
        )
        _set_current_worship_church(church_doc.id)
        _invalidate_worship_cache()
        flash(f"Joined {church.get('name') or church_doc.id}.", "success")
    except Exception as exc:
        app.logger.warning("worship_church_join error: %s", exc)
        flash("We couldn't join that church. Please try again.", "error")
    return redirect(url_for("worship"))


@app.route("/worship/church/switch", methods=["POST"])
@login_required
def worship_church_switch():
    church_id = _slugify_worship_token(request.form.get("church_id", ""))
    if not church_id:
        flash("Choose a church library.", "warning")
        return redirect(url_for("worship"))
    if church_id != _DEFAULT_WORSHIP_SCOPE and not _user_can_access_worship_church(church_id):
        flash("You are not a member of that church library.", "error")
        return redirect(url_for("worship"))
    _set_current_worship_church(church_id)
    _invalidate_worship_cache()
    flash("Switched worship library.", "success")
    return redirect(url_for("worship"))


def _slugify_worship_token(value: str) -> str:
    return re.sub(r"[^a-z0-9]+", "-", str(value or "").strip().lower()).strip("-")


def _worship_song_id_base(title: str, artist: str = "", version: str = "") -> str:
    title_slug = _slugify_worship_token(title)
    artist_slug = _slugify_worship_token(artist)
    version_slug = _slugify_worship_token(version)
    bits = [bit for bit in (title_slug, artist_slug, version_slug) if bit]
    return "-".join(bits) or "untitled-song"


def _make_unique_worship_song_id(title: str, artist: str = "", version: str = "") -> str:
    base = _worship_song_id_base(title, artist, version)
    candidate = base
    suffix = 2
    while get_worship_song(candidate):
        candidate = f"{base}-{suffix}"
        suffix += 1
    return candidate


def _worship_set_filename(selected_items: list[dict], date_label: str) -> str:
    title_slugs = []
    for item in selected_items[:2]:
        normalized = normalize_worship_song(item)
        slug = _slugify_worship_token(normalized.get("title", ""))
        if slug:
            title_slugs.append(slug)
    if title_slugs:
        return f"worship-{'-'.join(title_slugs)}.pptx"
    return f"worship-{date_label}.pptx"


def _clean_ai_json_response(raw: str) -> str:
    raw = str(raw or "").strip()
    raw = re.sub(r"^```[a-zA-Z0-9_-]*\n?", "", raw)
    raw = re.sub(r"\n?```$", "", raw).strip()
    return raw


_LYRICS_SITE_STOP_LINES = {
    "you may also like",
    "submit corrections",
    "submit lyrics",
    "soundtracks",
    "facebook",
    "contact us",
    "advertise here",
    "privacy policy",
    "cookie policy",
    "dmca policy",
}


def _is_lyrics_site_boilerplate_line(line: str) -> bool:
    clean = str(line or "").strip()
    lowered = clean.lower()
    if not clean:
        return False
    if lowered in _LYRICS_SITE_STOP_LINES:
        return True
    if lowered.startswith("azlyrics ") or lowered in {"azlyrics", "azlyrics.com"}:
        return True
    if lowered.startswith("writer(s):") or lowered.startswith("album:"):
        return True
    if lowered.startswith("play ") and ("apple music" in lowered or "spotify" in lowered):
        return True
    if lowered.startswith("lyrics licensed by") or lowered.startswith("copyright"):
        return True
    if re.fullmatch(r"(?:\.|\u2026|\s){3,}", clean):
        return True
    if re.fullmatch(r"\(?\s*(?:feat|featuring)\.?\s+.+\)?", clean, flags=re.I):
        return True
    if lowered.endswith(" lyrics") and len(clean.split()) <= 6:
        return True
    return False


def _clean_lyrics_site_paste(raw_text: str, title_hint: str = "", artist_hint: str = "") -> dict:
    """Strip common lyric-site chrome while preserving stanza breaks."""
    lines = [line.strip() for line in str(raw_text or "").replace("\r\n", "\n").replace("\r", "\n").split("\n")]
    inferred_title = str(title_hint or "").strip()
    inferred_artist = str(artist_hint or "").strip()

    for line in lines:
        if not inferred_title:
            match = re.match(r'^"(.+?)"\s+lyrics$', line, flags=re.I)
            if match:
                inferred_title = match.group(1).strip()
        if not inferred_artist:
            match = re.match(r"^(.+?)\s+Lyrics$", line)
            if match and not line.startswith('"') and "AZLyrics" not in line:
                inferred_artist = match.group(1).strip()

    start_idx = 0
    if inferred_title:
        quoted_title = f'"{inferred_title}"'.lower()
        for idx, line in enumerate(lines):
            if line.lower() == quoted_title:
                start_idx = idx + 1
                break
        else:
            for idx, line in enumerate(lines):
                if line.lower() == f'{quoted_title} lyrics':
                    start_idx = idx + 1
                    break

    cleaned: list[str] = []
    blank_pending = False
    for line in lines[start_idx:]:
        if not line:
            blank_pending = bool(cleaned)
            continue
        lowered = line.lower()
        if lowered in _LYRICS_SITE_STOP_LINES or lowered.startswith("you may also like") or lowered.startswith("writer(s):"):
            break
        if _is_lyrics_site_boilerplate_line(line):
            continue
        if blank_pending and cleaned and cleaned[-1] != "":
            cleaned.append("")
        cleaned.append(line)
        blank_pending = False

    while cleaned and cleaned[-1] == "":
        cleaned.pop()

    return {
        "title": inferred_title,
        "artist": inferred_artist,
        "lyrics": "\n".join(cleaned).strip(),
    }


def _split_clean_lyric_blocks(lyrics_text: str) -> list[list[str]]:
    blocks: list[list[str]] = []
    current: list[str] = []
    for raw_line in str(lyrics_text or "").replace("\r\n", "\n").replace("\r", "\n").split("\n"):
        line = raw_line.strip()
        if not line:
            if current:
                blocks.append(current)
                current = []
            continue
        if _is_lyrics_site_boilerplate_line(line):
            continue
        current.append(line)
    if current:
        blocks.append(current)
    return blocks


def _lyric_block_similarity(left: list[str], right: list[str]) -> float:
    left_set = {str(line or "").strip().lower() for line in left if str(line or "").strip()}
    right_set = {str(line or "").strip().lower() for line in right if str(line or "").strip()}
    if not left_set or not right_set:
        return 0.0
    return len(left_set & right_set) / max(len(left_set), len(right_set))


def _find_repeated_lyric_sequence(lines: list[str], min_len: int = 4, max_len: int = 10) -> tuple[int, int]:
    normalized = [str(line or "").strip().lower() for line in lines]
    max_len = min(max_len, len(normalized) // 2)
    for length in range(max_len, min_len - 1, -1):
        seen: dict[tuple[str, ...], int] = {}
        for idx in range(0, len(normalized) - length + 1):
            key = tuple(normalized[idx:idx + length])
            if key in seen and idx - seen[key] >= length:
                return seen[key], length
            seen.setdefault(key, idx)
    return -1, 0


def _lyric_line_equalish(left: str, right: str) -> bool:
    left_norm = re.sub(r"[^a-z0-9]+", " ", str(left or "").lower()).strip()
    right_norm = re.sub(r"[^a-z0-9]+", " ", str(right or "").lower()).strip()
    return bool(left_norm and right_norm and left_norm == right_norm)


def _extract_worship_section_label(line: str) -> tuple[str, bool]:
    text = str(line or "").strip()
    if not text:
        return "", False
    match = re.match(
        r"^\s*[\[(]?\s*((?:verse|v|chorus|ch|bridge|pre[-\s]?chorus|tag|intro|outro|ending|refrain)(?:\s*\d+)?)\s*(:?)\s*[\])]?\s*$",
        text,
        flags=re.I,
    )
    if not match:
        return "", False
    return _canonical_part_key(match.group(1)), bool(match.group(2))


def _parse_refrain_marker_worship_lyrics(
    lyrics_text: str,
    title: str = "",
    artist: str = "",
    version: str = "",
    key: str = "",
) -> dict | None:
    cleaned = _clean_lyrics_site_paste(lyrics_text, title, artist)
    lyric_body = cleaned.get("lyrics") or str(lyrics_text or "").strip()
    raw_lines = [
        raw_line.strip()
        for raw_line in lyric_body.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    ]
    lines = [line for line in raw_lines if line and not _is_lyrics_site_boilerplate_line(line)]
    marker_indexes: list[int] = []
    for idx, line in enumerate(lines):
        label, _has_colon = _extract_worship_section_label(line)
        if label == "chorus" and re.search(r"refrain", line, flags=re.I):
            marker_indexes.append(idx)

    if not marker_indexes or marker_indexes[0] <= 0:
        return None

    first_marker = marker_indexes[0]
    first_verse = lines[:first_marker]
    next_marker = marker_indexes[1] if len(marker_indexes) > 1 else len(lines)
    after_first_marker = lines[first_marker + 1:next_marker]
    if len(first_verse) < 2 or not after_first_marker:
        return None

    if len(marker_indexes) > 1 and len(after_first_marker) > len(first_verse):
        chorus_len = len(after_first_marker) - len(first_verse)
    else:
        chorus_len = len(after_first_marker)
    if chorus_len <= 0:
        return None

    chorus_lines = after_first_marker[:chorus_len]
    parts: dict[str, list[str]] = {"verse1": first_verse, "chorus": chorus_lines}
    arrangement: list[str] = ["verse1", "chorus"]
    verse_count = 2

    first_following_verse = after_first_marker[chorus_len:]
    if first_following_verse:
        part_name = f"verse{verse_count}"
        parts[part_name] = first_following_verse
        arrangement.append(part_name)
        verse_count += 1

    for marker_idx, next_idx in zip(marker_indexes[1:], marker_indexes[2:] + [len(lines)]):
        arrangement.append("chorus")
        following_verse = lines[marker_idx + 1:next_idx]
        if following_verse:
            part_name = f"verse{verse_count}"
            parts[part_name] = following_verse
            arrangement.append(part_name)
            verse_count += 1

    if len(parts) < 3 or arrangement.count("chorus") < 2:
        return None

    parsed_title = cleaned.get("title") or str(title or "").strip() or "Untitled Song"
    return {
        "id": re.sub(r"[^a-z0-9]+", "-", parsed_title.lower()).strip("-") or "untitled-song",
        "title": parsed_title,
        "artist": cleaned.get("artist") or str(artist or "").strip(),
        "version": str(version or "").strip(),
        "key": str(key or "").strip(),
        "type": "song",
        "parts": parts,
        "arrangement": arrangement,
    }


def _parse_continuous_worship_lyrics(
    lines: list[str],
    title: str = "",
    artist: str = "",
    version: str = "",
    key: str = "",
) -> dict | None:
    lines = [str(line or "").strip() for line in lines if str(line or "").strip()]
    if len(lines) < 16:
        return None
    chorus_start, chorus_len = _find_repeated_lyric_sequence(lines)
    if chorus_start < 0 or chorus_len < 4:
        return None

    chorus_lines = lines[chorus_start:chorus_start + chorus_len]
    parts: dict[str, list[str]] = {"chorus": chorus_lines}
    arrangement: list[str] = []
    verse_count = 1
    bridge_count = 1
    chorus_count = 0

    def is_chorus_at(index: int) -> bool:
        window = lines[index:index + chorus_len]
        if len(window) < max(4, chorus_len - 1):
            return False
        if not _lyric_line_equalish(window[0], chorus_lines[0]):
            return False
        return _lyric_block_similarity(window, chorus_lines) >= 0.6

    def add_block(block: list[str]) -> None:
        nonlocal verse_count, bridge_count, chorus_count
        if not block:
            return
        if chorus_count >= 2:
            part_name = "bridge" if bridge_count == 1 else f"bridge{bridge_count}"
            bridge_count += 1
            parts[part_name] = block
            arrangement.append(part_name)
            return
        for start in range(0, len(block), 4):
            chunk = block[start:start + 4]
            if not chunk:
                continue
            part_name = f"verse{verse_count}"
            verse_count += 1
            parts[part_name] = chunk
            arrangement.append(part_name)

    idx = 0
    while idx < len(lines):
        if is_chorus_at(idx):
            arrangement.append("chorus")
            chorus_count += 1
            idx += chorus_len
            continue
        next_chorus = -1
        for candidate in range(idx + 1, len(lines)):
            if is_chorus_at(candidate):
                next_chorus = candidate
                break
        end = next_chorus if next_chorus >= 0 else len(lines)
        add_block(lines[idx:end])
        idx = end

    if chorus_count < 2 or len(parts) < 3:
        return None
    return {
        "id": re.sub(r"[^a-z0-9]+", "-", str(title or "Untitled Song").lower()).strip("-") or "untitled-song",
        "title": str(title or "Untitled Song").strip() or "Untitled Song",
        "artist": str(artist or "").strip(),
        "version": str(version or "").strip(),
        "key": str(key or "").strip(),
        "type": "song",
        "parts": parts,
        "arrangement": arrangement,
    }


def _parse_labeled_worship_lyrics(lyrics_text: str, title: str = "", artist: str = "", version: str = "", key: str = "") -> dict | None:
    cleaned = _clean_lyrics_site_paste(lyrics_text, title, artist)
    lyric_body = cleaned.get("lyrics") or str(lyrics_text or "").strip()
    raw_sections: list[tuple[str, list[str]]] = []
    current_label = ""
    current_lines: list[str] = []

    def flush_current():
        nonlocal current_label, current_lines
        if current_label and current_lines:
            raw_sections.append((current_label, current_lines))
        current_label = ""
        current_lines = []

    for raw_line in lyric_body.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        section_label, _has_colon = _extract_worship_section_label(line)
        if section_label:
            flush_current()
            current_label = section_label
            continue
        if not current_label:
            continue
        current_lines.append(line)
    flush_current()

    if not raw_sections:
        return None

    parts: dict[str, list[str]] = {}
    arrangement: list[str] = []
    used_counts: dict[str, int] = {}
    content_to_part: dict[str, str] = {}
    for raw_label, lines in raw_sections:
        base_key = _canonical_part_key(raw_label)
        if not base_key:
            base_key = "verse"
        content_key = "\n".join(lines).lower()
        if content_key in content_to_part:
            part_name = content_to_part[content_key]
        else:
            used_counts[base_key] = used_counts.get(base_key, 0) + 1
            part_name = base_key if used_counts[base_key] == 1 else f"{base_key}{used_counts[base_key]}"
            parts[part_name] = lines
            content_to_part[content_key] = part_name
        arrangement.append(part_name)

    if not parts or not arrangement:
        return None

    parsed_title = cleaned.get("title") or str(title or "").strip() or "Untitled Song"
    return {
        "id": re.sub(r"[^a-z0-9]+", "-", parsed_title.lower()).strip("-") or "untitled-song",
        "title": parsed_title,
        "artist": cleaned.get("artist") or str(artist or "").strip(),
        "version": str(version or "").strip(),
        "key": str(key or "").strip(),
        "type": "song",
        "parts": parts,
        "arrangement": arrangement,
    }


def _fallback_parse_worship_lyrics(
    lyrics_text: str,
    title: str = "",
    artist: str = "",
    version: str = "",
    key: str = "",
) -> dict | None:
    refrain_markers = _parse_refrain_marker_worship_lyrics(lyrics_text, title, artist, version, key)
    if refrain_markers:
        return refrain_markers

    labeled = _parse_labeled_worship_lyrics(lyrics_text, title, artist, version, key)
    if labeled:
        return labeled

    cleaned = _clean_lyrics_site_paste(lyrics_text, title, artist)
    lyric_body = cleaned.get("lyrics") or str(lyrics_text or "").strip()
    blocks = _split_clean_lyric_blocks(lyric_body)
    if not blocks:
        return None
    if len(blocks) == 1 and len(blocks[0]) >= 16:
        continuous = _parse_continuous_worship_lyrics(
            blocks[0],
            cleaned.get("title") or title,
            cleaned.get("artist") or artist,
            version,
            key,
        )
        if continuous:
            return continuous
        chunked_parts = {
            f"verse{idx + 1}": blocks[0][start:start + 4]
            for idx, start in enumerate(range(0, len(blocks[0]), 4))
        }
        parsed_title = cleaned.get("title") or str(title or "").strip() or "Untitled Song"
        return {
            "id": re.sub(r"[^a-z0-9]+", "-", parsed_title.lower()).strip("-") or "untitled-song",
            "title": parsed_title,
            "artist": cleaned.get("artist") or str(artist or "").strip(),
            "version": str(version or "").strip(),
            "key": str(key or "").strip(),
            "type": "song",
            "parts": chunked_parts,
            "arrangement": list(chunked_parts.keys()),
        }

    # --- Pass 1: cluster near-duplicate blocks (similarity ≥ 0.65) ---
    # Canonical version of each cluster = the longest block seen so far for that cluster.
    canonicals: list[list[str]] = []
    block_cluster: list[int] = []
    for block in blocks:
        best_cidx, best_sim = -1, 0.0
        for cidx, canonical in enumerate(canonicals):
            sim = _lyric_block_similarity(block, canonical)
            if sim > best_sim:
                best_sim, best_cidx = sim, cidx
        if best_sim >= 0.65:
            block_cluster.append(best_cidx)
            if len(block) > len(canonicals[best_cidx]):
                canonicals[best_cidx] = block   # keep the most complete version
        else:
            block_cluster.append(len(canonicals))
            canonicals.append(list(block))

    # --- Pass 2: identify chorus = first cluster that appears more than once ---
    from collections import Counter
    cluster_counts = Counter(block_cluster)
    repeated_clusters = {cidx for cidx, cnt in cluster_counts.items() if cnt > 1}
    chorus_cluster = next(
        (block_cluster[i] for i in range(len(blocks)) if block_cluster[i] in repeated_clusters),
        -1,
    )

    # Find where the chorus first appears so we can distinguish bridge from verse
    chorus_first_idx = next(
        (i for i, cidx in enumerate(block_cluster) if cidx == chorus_cluster),
        len(blocks),
    )

    # --- Pass 3: assign part names ---
    cluster_to_part: dict[int, str] = {}
    parts: dict[str, list[str]] = {}
    arrangement: list[str] = []
    verse_count = 1
    bridge_count = 1

    for i, block in enumerate(blocks):
        cidx = block_cluster[i]
        if cidx in cluster_to_part:
            arrangement.append(cluster_to_part[cidx])
            continue

        # First occurrence of this cluster — assign a name
        if cidx == chorus_cluster:
            part_name = "chorus"
        elif cidx in repeated_clusters:
            # A repeated non-chorus section: bridge if it first appears after the chorus
            first_occ = next(j for j, c in enumerate(block_cluster) if c == cidx)
            if first_occ > chorus_first_idx:
                part_name = "bridge" if bridge_count == 1 else f"bridge{bridge_count}"
                bridge_count += 1
            else:
                part_name = f"verse{verse_count}"
                verse_count += 1
        else:
            part_name = f"verse{verse_count}"
            verse_count += 1

        cluster_to_part[cidx] = part_name
        parts[part_name] = canonicals[cidx]   # store the canonical (most complete) version
        arrangement.append(part_name)

    parsed_title = cleaned.get("title") or str(title or "").strip() or "Untitled Song"
    return {
        "id": re.sub(r"[^a-z0-9]+", "-", parsed_title.lower()).strip("-") or "untitled-song",
        "title": parsed_title,
        "artist": cleaned.get("artist") or str(artist or "").strip(),
        "version": str(version or "").strip(),
        "key": str(key or "").strip(),
        "type": "song",
        "parts": parts,
        "arrangement": arrangement,
    }


def _looks_like_line_exploded_worship_parse(parsed: dict) -> bool:
    parts = parsed.get("parts") if isinstance(parsed, dict) else None
    arrangement = parsed.get("arrangement") if isinstance(parsed, dict) else None
    if not isinstance(parts, dict) or not isinstance(arrangement, list):
        return False
    verse_like = []
    for name, lines in parts.items():
        canonical_name = _canonical_part_key(name)
        if re.match(r"^verse\d+$", canonical_name) and isinstance(lines, list):
            verse_like.append(name)
    if len(verse_like) < 10:
        return False
    one_line_parts = 0
    for name in verse_like:
        lines = parts.get(name) or []
        if isinstance(lines, list) and len([line for line in lines if str(line).strip()]) <= 1:
            one_line_parts += 1
    return one_line_parts / max(len(verse_like), 1) >= 0.75 and len(arrangement) >= 12


def _normalize_lyric_comparison_line(line: str) -> str:
    return re.sub(r"[^a-z0-9]+", " ", str(line or "").lower()).strip()


def _looks_under_arranged_worship_parse(parsed: dict, source_lyrics: str) -> bool:
    source_lines = {
        normalized
        for line in str(source_lyrics or "").splitlines()
        if (normalized := _normalize_lyric_comparison_line(line))
        and not _extract_worship_section_label(line)[0]
        and not _is_lyrics_site_boilerplate_line(line)
    }
    if len(source_lines) < 12 or not isinstance(parsed, dict):
        return False
    parts = parsed.get("parts")
    arrangement = parsed.get("arrangement")
    if not isinstance(parts, dict) or not isinstance(arrangement, list):
        return False
    if not arrangement:
        return True
    if len(arrangement) < 4:
        return True
    parsed_lines = {
        normalized
        for lines in parts.values()
        if isinstance(lines, list)
        for line in lines
        if (normalized := _normalize_lyric_comparison_line(line))
    }
    # Compare canonical content rather than arrangement-expanded content. Repeated
    # choruses must not compensate for unique source lines that the model dropped.
    coverage = len(source_lines & parsed_lines) / max(len(source_lines), 1)
    return coverage < 0.80


def _repair_line_exploded_worship_song(song: dict) -> dict | None:
    if not _looks_like_line_exploded_worship_parse(song):
        return None
    normalized = normalize_worship_song(song)
    old_parts = normalized.get("parts", {})
    old_arrangement = normalized.get("arrangement", [])

    ordered_lines: list[str] = []
    for part_name in old_arrangement:
        lines = old_parts.get(part_name, [])
        if isinstance(lines, list) and lines:
            ordered_lines.append(str(lines[0]).strip())
    ordered_lines = [line for line in ordered_lines if line]
    if len(ordered_lines) < 8:
        return None

    best_start = -1
    best_length = 0
    best_distance = 0
    max_len = min(10, max(4, len(ordered_lines) // 2))
    for length in range(max_len, 3, -1):
        seen: dict[tuple[str, ...], int] = {}
        for idx in range(0, len(ordered_lines) - length + 1):
            key = tuple(line.lower() for line in ordered_lines[idx:idx + length])
            if key in seen:
                distance = idx - seen[key]
                if distance > best_distance:
                    best_start = seen[key]
                    best_length = length
                    best_distance = distance
            else:
                seen[key] = idx
        if best_start >= 0:
            break

    if best_start < 0 or best_length < 4:
        return None

    chorus_lines = ordered_lines[best_start:best_start + best_length]
    rebuilt_parts: dict[str, list[str]] = {"chorus": chorus_lines}
    rebuilt_arrangement: list[str] = []
    verse_count = 1
    bridge_count = 1
    idx = 0

    def add_non_chorus_block(block: list[str], after_chorus: bool) -> None:
        nonlocal verse_count, bridge_count
        if not block:
            return
        chunk_size = 4 if len(block) <= 8 else 4
        for chunk_start in range(0, len(block), chunk_size):
            chunk = block[chunk_start:chunk_start + chunk_size]
            if not chunk:
                continue
            if after_chorus and len(chunk) >= 3:
                part_name = "bridge" if bridge_count == 1 else f"bridge{bridge_count}"
                bridge_count += 1
            else:
                part_name = f"verse{verse_count}"
                verse_count += 1
            rebuilt_parts[part_name] = chunk
            rebuilt_arrangement.append(part_name)

    while idx < len(ordered_lines):
        window = ordered_lines[idx:idx + best_length]
        if len(window) == best_length and _lyric_block_similarity(window, chorus_lines) >= 0.6:
            rebuilt_arrangement.append("chorus")
            idx += best_length
            continue
        next_chorus = -1
        for candidate in range(idx + 1, len(ordered_lines) - best_length + 1):
            candidate_window = ordered_lines[candidate:candidate + best_length]
            if _lyric_block_similarity(candidate_window, chorus_lines) >= 0.6:
                next_chorus = candidate
                break
        end = next_chorus if next_chorus >= 0 else len(ordered_lines)
        add_non_chorus_block(ordered_lines[idx:end], bool(rebuilt_arrangement))
        idx = end

    if len(rebuilt_parts) < 2 or "chorus" not in rebuilt_arrangement:
        return None

    repaired = dict(normalized)
    repaired["parts"] = rebuilt_parts
    repaired["arrangement"] = rebuilt_arrangement
    return normalize_worship_song(repaired)


class _ReadableHTMLTextParser(HTMLParser):
    def __init__(self):
        super().__init__()
        self._skip_depth = 0
        self._chunks: list[str] = []

    def handle_starttag(self, tag, attrs):
        tag = tag.lower()
        if tag in {"script", "style", "noscript", "svg", "canvas"}:
            self._skip_depth += 1
        if tag in {"br", "p", "div", "section", "article", "header", "footer", "li", "tr", "h1", "h2", "h3"}:
            self._chunks.append("\n")

    def handle_endtag(self, tag):
        tag = tag.lower()
        if tag in {"script", "style", "noscript", "svg", "canvas"} and self._skip_depth:
            self._skip_depth -= 1
        if tag in {"p", "div", "section", "article", "li", "tr", "h1", "h2", "h3"}:
            self._chunks.append("\n")

    def handle_data(self, data):
        if not self._skip_depth and data and data.strip():
            self._chunks.append(data.strip())
            self._chunks.append("\n")

    def text(self) -> str:
        lines: list[str] = []
        for raw_line in "\n".join(self._chunks).splitlines():
            line = re.sub(r"\s+", " ", raw_line).strip()
            if line:
                lines.append(line)
            elif lines and lines[-1] != "":
                lines.append("")
        while lines and lines[-1] == "":
            lines.pop()
        return "\n".join(lines)


def _extract_readable_text_from_html(html_text: str) -> str:
    parser = _ReadableHTMLTextParser()
    parser.feed(str(html_text or ""))
    return parser.text()


def _is_safe_worship_import_url(import_url: str) -> bool:
    parsed = urlparse(str(import_url or "").strip())
    if parsed.scheme not in {"http", "https"} or not parsed.hostname:
        return False
    try:
        addresses = socket.getaddrinfo(parsed.hostname, None)
    except OSError:
        return False
    for address in addresses:
        ip = ipaddress.ip_address(address[4][0])
        if ip.is_private or ip.is_loopback or ip.is_link_local or ip.is_multicast or ip.is_reserved:
            return False
    return True


class _SafeWorshipRedirectHandler(HTTPRedirectHandler):
    """Re-validate every redirect hop so a public URL can't bounce the importer
    to an internal address (SSRF via redirect)."""
    def redirect_request(self, req, fp, code, msg, headers, newurl):
        if not _is_safe_worship_import_url(newurl):
            raise ValueError("Blocked an unsafe redirect during song import.")
        return super().redirect_request(req, fp, code, msg, headers, newurl)


def _fetch_worship_import_text(import_url: str) -> str:
    import_url = str(import_url or "").strip()
    if not _is_safe_worship_import_url(import_url):
        raise ValueError("Enter a public http or https song page URL.")
    req = Request(
        import_url,
        headers={
            "User-Agent": "FaithSparksWorshipImporter/1.0 (+https://faithsparksprintables.com)",
            "Accept": "text/html,text/plain;q=0.9,*/*;q=0.6",
        },
    )
    opener = build_opener(_SafeWorshipRedirectHandler())
    with opener.open(req, timeout=10) as response:
        content_type = response.headers.get("Content-Type", "")
        raw_body = response.read(800_000)
    charset_match = re.search(r"charset=([^;]+)", content_type, flags=re.I)
    charset = charset_match.group(1).strip() if charset_match else "utf-8"
    body = raw_body.decode(charset, errors="replace")
    if "html" in content_type.lower() or "<html" in body[:500].lower():
        body = _extract_readable_text_from_html(body)
    return body.strip()


def _resolve_selected_worship_items(song_order: str, fallback_song_ids: list[str], *, prefer_direct: bool = False) -> list[dict]:
    if song_order:
        raw_ids = [s.strip() for s in song_order.split(",") if s.strip()]
    else:
        raw_ids = fallback_song_ids

    seen: set = set()
    ordered_ids = []
    for sid in raw_ids:
        if sid not in seen:
            seen.add(sid)
            ordered_ids.append(sid)

    # Resolve exact ids from the cached list first. Title-derived fallback is
    # allowed only when it is unambiguous, so duplicate song titles do not pick
    # whichever version happened to sort first.
    song_lookup: dict[str, dict] = {}
    title_lookup: dict[str, list[dict]] = {}
    try:
        for song in list_worship_songs():
            normalized = normalize_worship_song(song)
            song_id = normalized.get("id", "")
            if song_id:
                song_lookup[song_id] = song
            title = normalized.get("title", "")
            if title:
                title_lookup.setdefault(_slugify_worship_token(title), []).append(song)
    except Exception:
        pass

    selected_items = []
    for song_id in ordered_ids:
        title_matches = title_lookup.get(_slugify_worship_token(song_id), [])
        song = get_worship_song(song_id) if prefer_direct else None
        if not song:
            song = song_lookup.get(song_id)
        if not song and len(title_matches) == 1:
            song = title_matches[0]
        if not song:
            song = get_worship_song(song_id)
        if song:
            selected_items.append(song)
    return selected_items


def _build_worship_mobile_slides(selected_items: list[dict], notes: dict | None = None) -> list[dict]:
    notes = notes if isinstance(notes, dict) else {}
    slides: list[dict] = []
    for item in selected_items:
        normalized = normalize_worship_song(item)
        song_id = normalized.get("id", "")
        song_title = normalized.get("title", "Untitled")
        song_version = normalized.get("version", "")
        item_type = normalized.get("type", "song")
        song_bg = normalized.get("background")
        song_note = str(notes.get(song_id) or "").strip()
        slides.append(
            {
                "kind": "divider",
                "id": song_id,
                "title": song_title,
                "artist": normalized.get("artist", ""),
                "version": song_version,
                "key": normalized.get("key", ""),
                "type": item_type,
                "background": song_bg,
                "note": song_note,
            }
        )
        parts = normalized.get("parts", {}) or {}
        arrangement = normalized.get("arrangement", []) or []
        for part_name in arrangement:
            part_lines = parts.get(part_name, [])
            if not isinstance(part_lines, list):
                continue
            for chunk in chunk_lines(part_lines):
                slides.append(
                    {
                        "kind": "lyric",
                        "id": song_id,
                        "title": song_title,
                        "version": song_version,
                        "part": part_name,
                        "part_label": _worship_part_label(part_name),
                        "lines": chunk["lines"],
                        "font_size": chunk.get("font_size", 48),
                        "type": item_type,
                        "background": song_bg,
                        "note": song_note,
                    }
                )
    return slides


@app.route("/worship/mobile", methods=["GET"])
@login_required
def worship_mobile():
    _seed_worship_from_files()
    setlist_id = request.args.get("setlist_id", "").strip()
    setlist = _get_worship_setlist(setlist_id) if setlist_id else None
    if setlist:
        selected_items = _resolve_selected_worship_items("", setlist.get("songs", []))
        notes = setlist.get("notes", {})
    else:
        selected_items = _resolve_selected_worship_items(
            request.args.get("song_order", ""),
            request.args.getlist("song_ids"),
        )
        notes = {}
    if not selected_items:
        flash("Select at least one item to preview the mobile slides.", "warning")
        return redirect(url_for("worship"))
    slides = _build_worship_mobile_slides(selected_items, notes)
    song_order = ",".join(item.get("id", "") for item in selected_items if item.get("id"))
    return render_template(
        "worship_mobile.html",
        slides=slides,
        selected_items=selected_items,
        song_order=song_order,
        setlist=setlist,
    )


@app.route("/worship/mobile-qr.png", methods=["GET"])
@login_required
def worship_mobile_qr():
    qr_url = request.args.get("url", "").strip()
    parsed = urlparse(qr_url)
    host_url = urlparse(request.host_url)
    if not qr_url or parsed.netloc != host_url.netloc or parsed.path != url_for("worship_mobile"):
        return Response("Invalid QR target", status=400)
    try:
        import qrcode
    except Exception:
        app.logger.exception("qrcode package is not available")
        return Response("QR support is not installed", status=503)

    image = qrcode.make(qr_url)
    output = BytesIO()
    image.save(output, format="PNG")
    output.seek(0)
    return send_file(output, mimetype="image/png", max_age=300)


def _build_worship_lyric_sheet_pdf(selected_items: list[dict], mobile_url: str = "") -> BytesIO:
    """Render the selected songs into a formatted, printable lyric-sheet PDF."""
    import html as _html
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import inch
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.colors import HexColor
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, KeepTogether

    accent = HexColor("#0ea5a8")
    grey = HexColor("#64748b")
    dark = HexColor("#0f172a")
    title_style = ParagraphStyle("wTitle", fontName="Helvetica-Bold", fontSize=18, leading=22, textColor=dark, spaceAfter=2)
    sub_style = ParagraphStyle("wSub", fontName="Helvetica", fontSize=10, leading=13, textColor=grey, spaceAfter=8)
    label_style = ParagraphStyle("wLabel", fontName="Helvetica-Bold", fontSize=10.5, leading=13, textColor=accent, spaceBefore=9, spaceAfter=3)
    line_style = ParagraphStyle("wLine", fontName="Helvetica", fontSize=12, leading=16, textColor=dark)
    ref_style = ParagraphStyle("wRef", fontName="Helvetica-Oblique", fontSize=10.5, leading=14, textColor=grey, spaceAfter=2)
    footer_style = ParagraphStyle("wFoot", fontName="Helvetica", fontSize=8.5, leading=11, textColor=grey)

    def esc(value) -> str:
        return _html.escape(str(value or ""))

    story: list = []
    for idx, item in enumerate(selected_items):
        normalized = normalize_worship_song(item)
        header = [Paragraph(esc(normalized.get("title", "Untitled")), title_style)]
        sub_bits = []
        if normalized.get("artist"):
            sub_bits.append(normalized["artist"])
        if normalized.get("version"):
            sub_bits.append(normalized["version"])
        if normalized.get("key"):
            sub_bits.append(f"Key: {normalized['key']}")
        header.append(Paragraph(esc(" · ".join(sub_bits)), sub_style) if sub_bits else Spacer(1, 6))
        # Keep the title/subtitle with the first section so a song never orphans.
        story.append(KeepTogether(header))
        for block in build_lyric_sheet_blocks(normalized):
            story.append(Paragraph(esc(block["label"]).upper(), label_style))
            if block["reference_only"]:
                story.append(Paragraph("(repeat)", ref_style))
                continue
            for line in block["lines"]:
                story.append(Paragraph(esc(line), line_style))
        if idx != len(selected_items) - 1:
            story.append(Spacer(1, 22))

    if mobile_url:
        story.append(Spacer(1, 18))
        story.append(Paragraph("Mobile view: " + esc(mobile_url), footer_style))

    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=letter,
        leftMargin=0.85 * inch, rightMargin=0.85 * inch,
        topMargin=0.8 * inch, bottomMargin=0.7 * inch,
        title="Worship Lyric Sheet",
    )
    doc.build(story or [Spacer(1, 1)])
    buf.seek(0)
    return buf


@app.route("/worship/export/lyric-sheet", methods=["POST"])
@login_required
def worship_export_lyric_sheet():
    _seed_worship_from_files()
    selected_items = _resolve_selected_worship_items(
        request.form.get("song_order", ""),
        request.form.getlist("song_ids"),
        prefer_direct=True,
    )
    if not selected_items:
        flash("Select at least one item to export a lyric sheet.", "warning")
        return redirect(url_for("worship"))

    mobile_url = url_for(
        "worship_mobile",
        _external=True,
        song_order=",".join(item.get("id", "") for item in selected_items if item.get("id")),
    )
    today = datetime.now(timezone.utc).date().isoformat()
    pdf = _build_worship_lyric_sheet_pdf(selected_items, mobile_url)
    return send_file(
        pdf,
        as_attachment=True,
        download_name=f"worship_lyric_sheet_{today}.pdf",
        mimetype="application/pdf",
    )


@app.route("/worship/build", methods=["POST"])
@login_required
def worship_build():
    _seed_worship_from_files()
    selected_items = _resolve_selected_worship_items(
        request.form.get("song_order", ""),
        request.form.getlist("song_ids"),
    )

    if not selected_items:
        flash("Select at least one item to build a deck.", "warning")
        return redirect(url_for("worship"))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for item in selected_items:
        normalized = normalize_worship_song(item)
        item_type = normalized.get("type", "song")
        song_bg = normalized.get("background")
        artist_bits = [bit for bit in (normalized.get("artist", ""), normalized.get("version", "")) if bit]
        create_divider_slide(prs, normalized.get("title", ""), " | ".join(artist_bits), normalized.get("key", ""), item_type, song_bg)
        parts = normalized.get("parts", {}) or {}
        arrangement = normalized.get("arrangement", []) or []
        for part_name in arrangement:
            part_lines = parts.get(part_name, [])
            if not isinstance(part_lines, list):
                continue
            for slide in chunk_lines(part_lines):
                create_content_slide(
                    prs,
                    slide["lines"],
                    item_type,
                    song_bg,
                    font_size=slide.get("font_size", 48),
                    part_label=_worship_part_label(part_name),
                )

    today = datetime.now(timezone.utc).date().isoformat()
    _touch_worship_songs_last_used([item.get("id", "") for item in selected_items if item.get("id")], today)

    tmp = NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    tmp.close()

    @after_this_request
    def _cleanup(response):
        try:
            os.unlink(tmp.name)
        except Exception:
            pass
        return response

    return send_file(
        tmp.name,
        as_attachment=True,
        download_name=_worship_set_filename(selected_items, today),
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

@app.route("/worship/add", methods=["GET", "POST"])
@login_required
def worship_add():
    if request.method == "POST":
        # Overwrite confirmation branch
        if request.form.get("overwrite"):
            pending_token = session.get("pending_worship_token", "")
            try:
                pending_payload = _load_pending_worship_song(pending_token)
            except RuntimeError as exc:
                app.logger.warning("worship_add overwrite load failed: %s", exc)
                flash("The saved import could not be loaded right now. Please try again.", "error")
                return render_template("worship_add.html", conflict_song=None,
                                       backgrounds=list(_load_bg_config().keys())), 503
            if pending_payload:
                pending = pending_payload.get("song") or {}
                used_fallback = bool(pending_payload.get("used_fallback"))
                fallback_reason = str(pending_payload.get("fallback_reason") or "")
                try:
                    save_worship_song(pending)
                except RuntimeError as exc:
                    app.logger.warning("worship_add overwrite save failed: %s", exc)
                    flash("The song could not be overwritten right now. Your import is still available; please retry.", "error")
                    return render_template("worship_add.html", conflict_song=pending,
                                           backgrounds=list(_load_bg_config().keys())), 503
                _delete_pending_worship_song(pending_token)
                session.pop("pending_worship_token", None)
                if used_fallback:
                    flash(f"'{pending['title']}' overwritten. {fallback_reason or 'Faith Sparks auto-structured the sections.'} Please review.", "success")
                else:
                    flash(f"'{pending['title']}' overwritten.", "success")
                return redirect(url_for("worship"))
            flash("Nothing to overwrite.", "warning")
            return redirect(url_for("worship_add"))

        title = request.form.get("title", "").strip()
        artist = request.form.get("artist", "").strip()
        version = request.form.get("version", "").strip()
        key = request.form.get("key", "").strip()
        song_type = request.form.get("type", "song").strip() or "song"
        background = request.form.get("background", "").strip()

        part_names = request.form.getlist("part_name")
        part_lines_raw = request.form.getlist("part_lines")
        arrangement_raw = request.form.get("arrangement", "")

        if not title:
            flash("Title is required.", "warning")
            return redirect(url_for("worship_add"))

        parts = {}
        for name, lines_text in zip(part_names, part_lines_raw):
            name = _canonical_part_key(name)
            if not name:
                continue
            lines = [l.strip() for l in lines_text.splitlines() if l.strip()]
            if lines:
                existing = parts.setdefault(name, [])
                existing.extend(line for line in lines if line not in existing)

        arrangement = [a.strip() for a in arrangement_raw.split(",") if a.strip()]

        song_id = _worship_song_id_base(title, artist, version) or "untitled-song"
        song = normalize_worship_song({
            "id": song_id,
            "title": title,
            "artist": artist,
            "version": version,
            "key": key,
            "type": song_type,
            "background": background,
            "parts": parts,
            "arrangement": arrangement,
        })

        if not song["parts"]:
            flash("Add at least one song part with lyrics.", "warning")
            return redirect(url_for("worship_add"))

        # Re-adding an existing song routes through the overwrite/conflict screen
        # rather than silently creating title-2, title-3, ...
        try:
            existing_song = get_worship_song(song_id)
            if existing_song:
                session["pending_worship_token"] = _store_pending_worship_song(song, False, "")
                return redirect(url_for("worship_add", conflict=song_id))
            save_worship_song(song)
        except RuntimeError as exc:
            app.logger.warning("worship_add persistence failed: %s", exc)
            flash("The song library is temporarily unavailable. Please try again.", "error")
            return render_template("worship_add.html", conflict_song=None,
                                   backgrounds=list(_load_bg_config().keys())), 503
        flash(f"'{title}' saved.", "success")
        return redirect(url_for("worship"))

    conflict_song = None
    conflict_id = request.args.get("conflict", "")
    if conflict_id:
        try:
            conflict_song = get_worship_song(conflict_id)
        except RuntimeError as exc:
            app.logger.warning("worship_add conflict lookup failed: %s", exc)
            flash("The song library is temporarily unavailable. Please try again.", "error")
            return render_template("worship_add.html", conflict_song=None,
                                   backgrounds=list(_load_bg_config().keys())), 503
    return render_template("worship_add.html", conflict_song=conflict_song,
                            backgrounds=list(_load_bg_config().keys()))


def _parse_worship_lyrics_claude(prompt: str, api_key: str) -> dict:
    """Parse lyrics with Claude. Returns a parsed dict or raises."""
    import anthropic
    model = os.environ.get("WORSHIP_PARSE_MODEL", "claude-sonnet-4-6")
    client = anthropic.Anthropic(api_key=api_key)
    msg = client.messages.create(
        model=model,
        max_tokens=4096,
        temperature=0,
        system="You are a worship song librarian. Output ONLY valid JSON — no prose, no markdown fences.",
        # Newer Claude models require the conversation to end with a user turn
        # and reject the older assistant-prefill JSON technique.
        messages=[{"role": "user", "content": prompt}],
    )
    text = "".join(
        block.text for block in msg.content if getattr(block, "type", None) == "text"
    )
    raw_json = _clean_ai_json_response(text)
    app.logger.info("worship_add_parse: Claude (%s) returned %d chars", model, len(raw_json))
    return json.loads(raw_json)


def _parse_worship_lyrics_openai(prompt: str, api_key: str) -> dict:
    """Parse lyrics with OpenAI gpt-4o, with a gpt-4o-mini repair pass. Returns a parsed dict or raises."""
    from openai import OpenAI
    client = OpenAI(api_key=api_key)
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are a worship song librarian. Output only valid JSON."},
            {"role": "user", "content": prompt},
        ],
        temperature=0.3,
        response_format={"type": "json_object"},
        max_tokens=4000,
    )
    if response.choices[0].finish_reason == "length":
        app.logger.warning("worship_add_parse: gpt-4o hit max_tokens, response truncated")
    raw_json = _clean_ai_json_response(response.choices[0].message.content)
    app.logger.info("worship_add_parse: gpt-4o returned %d chars", len(raw_json))
    try:
        return json.loads(raw_json)
    except json.JSONDecodeError:
        repair_prompt = (
            "Repair this malformed response into valid JSON only.\n\n"
            "Return the same worship song schema with id, title, artist, version, "
            "key, type, parts, and arrangement.\n"
            "Do not add markdown or explanation.\n\n"
            f"Malformed response:\n{raw_json}\n"
        )
        repaired = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": repair_prompt}],
            temperature=0,
            response_format={"type": "json_object"},
        )
        return json.loads(_clean_ai_json_response(repaired.choices[0].message.content))


@app.route("/worship/add/parse", methods=["POST"])
@login_required
def worship_add_parse():
    if request.content_length and request.content_length > 100_000:
        flash("That import is too large. Please paste only the song lyrics.", "warning")
        return render_template("worship_add.html", conflict_song=None,
                               backgrounds=list(_load_bg_config().keys())), 413
    raw_lyrics = request.form.get("raw_lyrics", "").strip()
    import_url = request.form.get("import_url", "").strip()
    title = request.form.get("title", "").strip()
    artist = request.form.get("artist", "").strip()
    version = request.form.get("version", "").strip()
    key = request.form.get("key", "").strip()
    submitted_title, submitted_artist = title, artist
    submitted_version, submitted_key = version, key

    user_key = str(session.get("user_email") or get_client_ip())
    user_rate = check_rate_limit("worship_import:user", user_key, limit=30, window_seconds=3600)
    ip_rate = check_rate_limit("worship_import:ip", get_client_ip(), limit=60, window_seconds=3600)
    if not user_rate.allowed or not ip_rate.allowed:
        retry_after = max(user_rate.retry_after, ip_rate.retry_after)
        flash("Too many song imports were requested. Please wait a little and try again.", "warning")
        response = app.make_response((render_template(
            "worship_add.html", conflict_song=None,
            backgrounds=list(_load_bg_config().keys())
        ), 429))
        response.headers["Retry-After"] = str(retry_after)
        return response

    if len(raw_lyrics) > 60_000:
        flash("That paste is too large. Please include only the song page or lyrics.", "warning")
        return redirect(url_for("worship_add"))

    if import_url:
        try:
            fetched_text = _fetch_worship_import_text(import_url)
        except Exception as e:
            app.logger.warning("worship_add_parse: import failed: %s", e, exc_info=True)
            flash("Could not read that song link. Please paste the lyrics directly and try again.", "error")
            return redirect(url_for("worship_add"))
        raw_lyrics = "\n\n".join(part for part in [raw_lyrics, fetched_text] if part)

    if len(raw_lyrics) > 60_000 or len(raw_lyrics.splitlines()) > 1_200:
        flash("That song page is too large to import. Please paste only the lyrics.", "warning")
        return redirect(url_for("worship_add"))

    if not raw_lyrics:
        flash("Paste lyrics or enter a song page link to parse.", "warning")
        return redirect(url_for("worship_add"))

    cleaned_paste = _clean_lyrics_site_paste(raw_lyrics, title, artist)
    parse_lyrics = cleaned_paste.get("lyrics") or raw_lyrics

    # AZLyrics (and similar sites) wrap each lyric line in its own <div>, so a
    # copy-paste produces a blank line between EVERY lyric line, not just between
    # stanzas. Detect this: if ≥40% of the lines are blank, we have the
    # every-line-blank format. Strip all blanks so gpt-4o sees consecutive text.
    all_lines = parse_lyrics.splitlines()
    blank_count = sum(1 for ln in all_lines if not ln.strip())
    if all_lines and blank_count / len(all_lines) >= 0.40:
        parse_lyrics = "\n".join(ln for ln in all_lines if ln.strip())
        app.logger.info("worship_add_parse: stripped every-line-blanks (%d blanks removed)", blank_count)

    title = title or cleaned_paste.get("title", "")
    artist = artist or cleaned_paste.get("artist", "")
    app.logger.info("worship_add_parse: cleaned lyrics (%d chars, %d lines)",
                    len(parse_lyrics), parse_lyrics.count("\n"))

    anthropic_key = os.environ.get("ANTHROPIC_API_KEY")
    openai_key = os.environ.get("OPENAI_API_KEY")

    prompt = f"""You are a worship song librarian. Parse the following song lyrics into structured JSON for a slide builder.

Title: {title or '(infer from lyrics)'}
Artist: {artist or '(unknown)'}
Version/arrangement note: {version or '(unknown)'}
Key: {key or '(unknown)'}

Lyrics:
{parse_lyrics}

Return ONLY valid JSON — no markdown fences, no explanation — in this exact format:
{{
  "id": "<slugified-title>",
  "title": "<title>",
  "artist": "<artist or empty string>",
  "version": "<version or arrangement note, or empty string>",
  "key": "<key or empty string>",
  "type": "song",
  "parts": {{
    "verse1": ["line 1", "line 2", "line 3", "line 4"],
    "chorus": ["line 1", "line 2", "line 3", "line 4"],
    "bridge": ["line 1", "line 2", "line 3", "line 4"]
  }},
  "arrangement": ["verse1", "chorus", "verse2", "chorus", "bridge", "bridge", "chorus"]
}}

CRITICAL RULES — read carefully before outputting:

DEDUPLICATION (most important rule):
- Each unique lyric section is stored ONCE in parts, then referenced as many times as needed in arrangement.
- If the chorus appears 3 times in the song, parts has ONE "chorus" entry and arrangement lists "chorus" three times.
- If a bridge repeats 2-3 times in a row, store it ONCE in parts and repeat the key in arrangement. Do NOT create bridge2, bridge3, bridge4 for repetitions of the same content.
- Only create numbered variants (chorus2, bridge2) when the lyric content is GENUINELY DIFFERENT — different words, not just a repeat.

STANZA BOUNDARIES:
- Blank lines, when present, separate stanzas — treat them as hard boundaries.
- The lyrics may arrive with NO blank line separators (copy-paste from web pages often strips whitespace). In that case, use musical structure and repeated content to identify section breaks: verses typically have 4 lines, choruses are the repeated block, bridges come after the second chorus.
- Never split a stanza mid-thought across multiple parts.

COMPLETENESS:
- Every lyric line in the source must appear in exactly one part. Do not drop lines or truncate stanzas.
- Count the lines. If verse 2 has 4 lines in the source, it must have 4 lines in the output.

NEAR-DUPLICATE SECTIONS:
- If the same section repeats with only minor lyric variations (e.g., an extra "Oh," at the start, or a slightly different last line), treat them as the SAME part. Use the most complete version as the canonical text.

OTHER RULES:
- id: title lowercased, spaces/special chars replaced by hyphens, no leading/trailing hyphens
- parts keys: use verse1, verse2, verse3, chorus, chorus2, bridge, pre_chorus, tag, outro, intro as appropriate
- each part value is an array of individual lyric lines (no blank strings)
- preserve apostrophes and contractions exactly as written
- ignore ALL website boilerplate: recommendations, ads, copyright, navigation, "You May Also Like", writer credits, album info
- type is always "song"
"""

    used_fallback_parser = False
    fallback_reason = ""

    # Provider order: Claude first (strongest at structured lyric extraction),
    # then OpenAI, then the local heuristic parser.
    parsed = None
    if anthropic_key:
        try:
            parsed = _parse_worship_lyrics_claude(prompt, anthropic_key)
        except Exception as e:
            app.logger.warning("worship_add_parse: Claude parse failed (%s), trying next provider: %s",
                               type(e).__name__, e, exc_info=True)

    if parsed is None and openai_key:
        try:
            parsed = _parse_worship_lyrics_openai(prompt, openai_key)
        except Exception as e:
            app.logger.error("worship_add_parse: OpenAI parse failed: %s", e, exc_info=True)

    if parsed is None:
        parsed = _fallback_parse_worship_lyrics(parse_lyrics, title, artist, version, key)
        if not parsed:
            if not (anthropic_key or openai_key):
                flash("No AI key is configured and the local parser could not find song sections.", "error")
            else:
                flash("We couldn't structure those lyrics yet. Please check the text and try again.", "error")
            return redirect(url_for("worship_add"))
        fallback_reason = "AI parsing was unavailable, so Faith Sparks auto-structured the sections."
        flash(fallback_reason, "warning")
        used_fallback_parser = True

    # Validate on the NORMALIZED song, not the raw model output: normalization
    # canonicalizes part keys and aligns the arrangement to them, so a model that
    # returns e.g. arrangement ["Chorus"] with parts {"chorus": [...]} (case/spacing
    # mismatch) is no longer falsely flagged "incomplete" and bounced to the fallback.
    song = normalize_worship_song(parsed) if isinstance(parsed, dict) else {"parts": {}, "arrangement": []}
    exploded = _looks_like_line_exploded_worship_parse(song)
    under_arranged = _looks_under_arranged_worship_parse(song, parse_lyrics)
    app.logger.info("worship_add_parse: parts=%s arrangement_len=%d exploded=%s under_arranged=%s",
                    list(song.get("parts", {}).keys()),
                    len(song.get("arrangement", [])),
                    exploded,
                    under_arranged)
    if (
        not song.get("parts")
        or not song.get("arrangement")   # empty arrangement = incomplete parse
        or exploded
        or under_arranged
    ):
        fallback = _fallback_parse_worship_lyrics(parse_lyrics, title, artist, version, key)
        if not fallback:
            flash("AI response was missing required fields (parts/arrangement).", "error")
            return redirect(url_for("worship_add"))
        song = normalize_worship_song(fallback)
        fallback_reason = "AI response was incomplete, so Faith Sparks auto-structured the sections."
        used_fallback_parser = True
        if (
            not song.get("parts")
            or not song.get("arrangement")
            or _looks_under_arranged_worship_parse(song, parse_lyrics)
        ):
            flash("We could not structure all of the lyrics reliably. Please add section labels or use manual entry.", "error")
            return redirect(url_for("worship_add"))

    # The form explicitly promises that supplied metadata overrides inference.
    for field, supplied_value in (
        ("title", submitted_title),
        ("artist", submitted_artist),
        ("version", submitted_version),
        ("key", submitted_key),
    ):
        if supplied_value:
            song[field] = supplied_value
    if not song.get("title"):
        flash("We could not determine the song title. Enter a title and try again.", "warning")
        return redirect(url_for("worship_add"))
    # Use the base ID (not a unique-suffixed one) so re-importing the same song
    # triggers the conflict/overwrite flow instead of creating gratitude-...-2, -3, etc.
    song["id"] = _worship_song_id_base(song.get("title", ""), song.get("artist", ""), song.get("version", "")) or "untitled-song"

    try:
        if get_worship_song(song["id"]):
            session["pending_worship_token"] = _store_pending_worship_song(
                song, used_fallback_parser, fallback_reason
            )
            return redirect(url_for("worship_add", conflict=song["id"]))
        save_worship_song(song)
    except RuntimeError as exc:
        app.logger.warning("worship_add_parse persistence failed: %s", exc)
        flash("The song was parsed, but the library is temporarily unavailable. Please retry the import.", "error")
        return render_template("worship_add.html", conflict_song=None,
                               backgrounds=list(_load_bg_config().keys())), 503
    if used_fallback_parser:
        flash(f"'{song.get('title', song['id'])}' saved. {fallback_reason or 'Faith Sparks auto-structured the sections.'} Please review.", "success")
    else:
        flash(f"'{song.get('title', song['id'])}' parsed and saved.", "success")
    return redirect(url_for("worship"))


@app.route("/worship/preview-slides", methods=["POST"])
@login_required
def worship_preview_slides():
    """Return slide chunks for a block of lyric lines as JSON (used by live preview UI)."""
    lines_text = request.json.get("lines", "") if request.is_json else request.form.get("lines", "")
    lines = [l for l in str(lines_text).splitlines() if l.strip()]
    slides = chunk_lines(lines)
    return jsonify({"slides": slides})


@app.route("/worship/delete", methods=["POST"])
@login_required
def worship_delete():
    song_id = request.form.get("song_id", "").strip()
    if not song_id or ".." in song_id or "/" in song_id or "\\" in song_id:
        if _worship_wants_json_response():
            return jsonify({"ok": False, "error": "Invalid song id"}), 400
        flash("Invalid song id.", "warning")
        return redirect(url_for("worship"))

    if delete_worship_song(song_id):
        setlists_changed = _remove_song_from_worship_setlists(song_id)
        if _worship_wants_json_response():
            return jsonify({"ok": True, "song_id": song_id, "setlists_changed": setlists_changed})
        flash("Song deleted.", "success")
    else:
        if _worship_wants_json_response():
            return jsonify({"ok": False, "error": "Song not found"}), 404
        flash("Song not found.", "warning")
    return redirect(url_for("worship"))


@app.route("/worship/library/reset", methods=["POST"])
@login_required
def worship_library_reset():
    confirmation = request.form.get("confirmation", "").strip().upper()
    if confirmation != "DELETE":
        if _worship_wants_json_response():
            return jsonify({"ok": False, "error": "Type DELETE to reset this worship library."}), 400
        flash("Type DELETE to reset this worship library.", "warning")
        return redirect(url_for("worship"))
    deleted = _delete_all_worship_library_data()
    if _worship_wants_json_response():
        return jsonify({"ok": True, **deleted})
    flash(f"Reset worship library: deleted {deleted['songs']} songs and {deleted['setlists']} saved setlists.", "success")
    return redirect(url_for("worship"))


@app.route("/worship/duplicate/<song_id>", methods=["POST"])
@login_required
def worship_duplicate(song_id):
    if ".." in song_id or "/" in song_id or "\\" in song_id:
        flash("Invalid song id.", "warning")
        return redirect(url_for("worship"))
    song = get_worship_song(song_id)
    if not song:
        flash("Song not found.", "warning")
        return redirect(url_for("worship"))
    duplicate = normalize_worship_song(song)
    base_version = duplicate.get("version", "")
    duplicate["version"] = f"{base_version} Copy".strip() if base_version else "Copy"
    duplicate.pop("last_used", None)
    duplicate["id"] = _make_unique_worship_song_id(
        duplicate.get("title", ""),
        duplicate.get("artist", ""),
        duplicate.get("version", ""),
    )
    save_worship_song(duplicate)
    flash(f"'{duplicate.get('title', duplicate['id'])}' duplicated.", "success")
    return redirect(url_for("worship_edit", song_id=duplicate["id"]))


@app.route("/worship/edit/<song_id>", methods=["GET", "POST"])
@login_required
def worship_edit(song_id):
    if ".." in song_id or "/" in song_id or "\\" in song_id:
        flash("Invalid song id.", "warning")
        return redirect(url_for("worship"))

    song = get_worship_song(song_id)
    if not song:
        flash("Song not found.", "warning")
        return redirect(url_for("worship"))

    if request.method == "GET":
        pass  # auto-repair disabled: it produced worse results than the incomplete parse

    if request.method == "POST":
        title = request.form.get("title", "").strip()
        artist = request.form.get("artist", "").strip()
        version = request.form.get("version", "").strip()
        key = request.form.get("key", "").strip()
        song_type = request.form.get("type", "song").strip() or "song"
        background = request.form.get("background", "").strip()

        part_names = request.form.getlist("part_name")
        part_lines_raw = request.form.getlist("part_lines")
        arrangement_raw = request.form.get("arrangement", "")

        if not title:
            flash("Title is required.", "warning")
            return redirect(url_for("worship_edit", song_id=song_id))

        parts = {}
        for name, lines_text in zip(part_names, part_lines_raw):
            name = name.strip()
            if not name:
                continue
            lines = [l.strip() for l in lines_text.splitlines() if l.strip()]
            if lines:
                parts[name] = lines

        arrangement = [a.strip() for a in arrangement_raw.split(",") if a.strip()]

        song.update({
            "title": title,
            "artist": artist,
            "version": version,
            "key": key,
            "type": song_type,
            "background": background,
            "parts": parts,
            "arrangement": arrangement,
        })

        save_worship_song(song)
        flash(f"'{title}' updated.", "success")
        return redirect(url_for("worship"))

    return render_template("worship_edit.html", song=song,
                            backgrounds=list(_load_bg_config().keys()))


@app.route("/worship/setlist/save", methods=["POST"])
@login_required
def worship_setlist_save():
    song_ids = request.form.getlist("song_ids")
    if not song_ids:
        return jsonify({"ok": False, "error": "No songs"}), 400
    setlist_name = request.form.get("setlist_name", "").strip()
    existing_id = request.form.get("setlist_id", "").strip()

    notes_json = request.form.get("notes_json", "{}")
    try:
        notes = json.loads(notes_json)
        if not isinstance(notes, dict):
            notes = {}
    except (json.JSONDecodeError, ValueError):
        notes = {}

    existing = _get_worship_setlist(existing_id) if existing_id else None
    today = datetime.now().strftime("%Y-%m-%d")
    date_label = existing.get("date") if existing else today
    setlist_id = existing_id if existing else _worship_setlist_id(today, setlist_name)
    created_at = existing.get("created_at") if existing else ""
    created_by = existing.get("created_by") if existing else ""
    now = _worship_timestamp()
    data = {
        "id": setlist_id,
        "date": date_label,
        "name": setlist_name,
        "songs": song_ids,
        "notes": notes,
        "worship_scope": _current_worship_scope(),
        "created_at": created_at or now,
        "updated_at": now,
        "created_by": created_by or session.get("user_email", ""),
        "updated_by": session.get("user_email", ""),
        "song_count": len(song_ids),
    }
    existed = _persist_worship_setlist(data, existing_id if existing else "")
    _touch_worship_songs_last_used(song_ids, today)

    return jsonify({"ok": True, **_normalize_worship_setlist(data), "updated": existed})


@app.route("/worship/setlist/delete", methods=["POST"])
@login_required
def worship_setlist_delete():
    setlist_id = request.form.get("setlist_id", "").strip() or request.form.get("date", "").strip()
    if not _valid_worship_setlist_id(setlist_id):
        return jsonify({"ok": False, "error": "Invalid setlist"}), 400
    _delete_worship_setlist(setlist_id)
    return jsonify({"ok": True})


@app.route("/worship/setlist/rename", methods=["POST"])
@login_required
def worship_setlist_rename():
    setlist_id = request.form.get("setlist_id", "").strip()
    new_name = request.form.get("setlist_name", "").strip()
    if not _valid_worship_setlist_id(setlist_id) or not new_name:
        return jsonify({"ok": False, "error": "Invalid setlist"}), 400
    existing = _get_worship_setlist(setlist_id)
    if not existing:
        return jsonify({"ok": False, "error": "Setlist not found"}), 404
    new_id = _make_unique_worship_setlist_id(existing["date"], new_name, setlist_id)
    existing.update(
        {
            "id": new_id,
            "name": new_name,
            "updated_at": _worship_timestamp(),
            "updated_by": session.get("user_email", ""),
            "song_count": len(existing.get("songs", [])),
            "worship_scope": _current_worship_scope(),
        }
    )
    _persist_worship_setlist(existing, setlist_id)
    return jsonify({"ok": True, **_normalize_worship_setlist(existing), "updated": True})


@app.route("/worship/setlist/duplicate", methods=["POST"])
@login_required
def worship_setlist_duplicate():
    setlist_id = request.form.get("setlist_id", "").strip()
    duplicate_name = request.form.get("setlist_name", "").strip()
    if not _valid_worship_setlist_id(setlist_id):
        return jsonify({"ok": False, "error": "Invalid setlist"}), 400
    existing = _get_worship_setlist(setlist_id)
    if not existing:
        return jsonify({"ok": False, "error": "Setlist not found"}), 404
    today = datetime.now().strftime("%Y-%m-%d")
    base_name = existing.get("name") or existing.get("date") or setlist_id
    duplicate_name = duplicate_name or f"Copy of {base_name}"
    new_id = _make_unique_worship_setlist_id(today, duplicate_name)
    now = _worship_timestamp()
    duplicate = {
        **existing,
        "id": new_id,
        "date": today,
        "name": duplicate_name,
        "created_at": now,
        "updated_at": now,
        "created_by": session.get("user_email", ""),
        "updated_by": session.get("user_email", ""),
        "song_count": len(existing.get("songs", [])),
        "worship_scope": _current_worship_scope(),
    }
    _persist_worship_setlist(duplicate)
    return jsonify({"ok": True, **_normalize_worship_setlist(duplicate), "updated": False})


## about + healthz moved to public blueprint

@app.route("/generate", methods=["GET", "POST"])
@login_required
def generate():
    from faithsparks.views.worksheets import generate as _impl
    return _impl()


@app.route("/illustrate", methods=["GET", "POST"])
@login_required
def illustrate():
    """Legacy route retained only for redirecting to the generate page."""
    flash("Coloring now lives on the Generate page. Browse the new library there.", "info")
    target = url_for("generate") + "#coloring"
    if request.method == "POST":
        return jsonify({"redirect": target, "message": "Coloring moved to /generate"}), 410
    return redirect(target)

@app.route("/delete/<filename>", methods=["POST"])
@login_required
def delete_worksheet(filename):
    from faithsparks.views.worksheets import delete_worksheet as _impl
    return _impl(filename)

@app.route("/delete_bulk", methods=["POST"])
@login_required
def delete_bulk():
    from faithsparks.views.worksheets import delete_bulk as _impl
    return _impl()
    

@app.route("/prints")
@login_required
def prints():
    from faithsparks.views.worksheets import history as _impl
    return _impl()

@app.route("/history")
def history():
    return redirect(url_for("prints"))
@app.route("/download/<filename>")
@login_required
def download_file(filename):
    from faithsparks.views.worksheets import download_file as _impl
    return _impl(filename)


@app.route("/coloring/image/<path:filename>")
@login_required
def coloring_image(filename):
    safe_name = os.path.basename(filename)
    if not safe_name or os.path.splitext(safe_name)[1].lower() != ".png":
        abort(404)
    if db:
        try:
            docs = (
                db.collection("worksheets")
                .where(filter=firestore.FieldFilter("email", "==", session.get("user_email")))
                .where(filter=firestore.FieldFilter("imageFilename", "==", safe_name))
                .limit(1)
                .stream()
            )
            if not next(docs, None):
                abort(404)
        except Exception:
            abort(404)
    local_path = os.path.join("worksheets", safe_name)
    if os.path.exists(local_path):
        return send_file(local_path, mimetype="image/png", conditional=True)
    remote = signed_url_for_path(f"worksheets/{safe_name}")
    if remote:
        return redirect(remote)
    flash("Image not found. It may have been removed.", "error")
    return redirect(url_for("prints"))

@app.route('/thumb/<path:filename>')
@login_required
def thumb(filename):
    from faithsparks.views.worksheets import thumb as _impl
    return _impl(filename)

# --- Admin utilities ---
def is_admin_email(email: str) -> bool:
    allow = os.getenv('ADMIN_EMAILS', '')
    if not allow:
        return False
    allowed = [e.strip().lower() for e in allow.split(',') if e.strip()]
    return (email or '').lower() in allowed

@app.route('/admin/seed_collections')
@admin_required
def admin_seed_collections():
    from faithsparks.views.admin.collections import admin_seed_collections as _impl
    return _impl()

@app.context_processor
def inject_helpers():
    fb_purchase = session.pop('fb_purchase', None)
    fb_user_match = None
    try:
        email = session.get('user_email')
        is_pro = False
        theme_name, theme_vars = get_theme_selection()
        # Resolve themed logo and favicon
        logo_url = url_for('static', filename='faith_sparks_logo.png')
        favicon_url = url_for('static', filename='favicon.ico')
        conf = _get_cached_config('app') or {}
        logos = (conf or {}).get('logos') or {}
        if isinstance(logos, dict):
            logo_url = logos.get(theme_name) or logos.get('default') or logo_url
        favs = (conf or {}).get('favicons') or {}
        if isinstance(favs, dict):
            favicon_url = favs.get(theme_name) or favs.get('default') or favicon_url
        site_content = _get_cached_config('content') or {}
        usage_nav = None
        path = request.path or ""
        if db and email:
            # On ordinary (non-usage) pages, serve the slow-changing is_pro flag
            # from a short-lived session cache so navigation doesn't read Firestore
            # on every render. Usage pages always read (they need live usage).
            fetch_usage = _should_fetch_usage(path)
            now_ts = time.time()
            _uc = session.get('_uc') if isinstance(session.get('_uc'), dict) else None
            user_doc = None
            if fetch_usage or not (_uc and (now_ts - _uc.get('ts', 0)) < _USER_FLAGS_TTL):
                # Read the user doc once (shared with the view via the request-scoped
                # cache) and refresh the session pro/plan cache.
                user_doc = get_user_doc(email)
                # Gift expiry downgrade (mirrors usage._get_user_plan side effect).
                try:
                    exp = user_doc.get('giftExpiresAt')
                    if exp and hasattr(exp, 'timestamp') and datetime.now(timezone.utc) > exp:
                        db.collection('users').document(email).set(
                            {'plan': 'free', 'isPro': False, 'giftExpiresAt': None}, merge=True)
                        user_doc['plan'] = 'free'
                        user_doc['isPro'] = False
                        user_doc['giftExpiresAt'] = None
                except Exception:
                    pass
                is_pro = bool(user_doc.get('isPro'))
                session['_uc'] = {
                    'isPro': is_pro,
                    'plan': user_doc.get('plan') or ('family' if is_pro else 'free'),
                    'ts': now_ts,
                }
            else:
                is_pro = bool(_uc.get('isPro'))

            # usage chip (sometimes expensive)
            if fetch_usage and user_doc is not None:
                try:
                    plan = user_doc.get('plan') or ('family' if user_doc.get('isPro') else 'free')
                    m_lim, _ = _quota_for_plan(plan)
                    usage = user_doc.get('usage') or {}
                    used_m = int((usage.get('months') or {}).get(_month_key()) or 0)
                    if m_lim is not None:
                        try:
                            used_val = int(used_m)
                        except Exception:
                            used_val = 0
                        try:
                            limit_val = int(m_lim)
                        except Exception:
                            limit_val = 0
                        remaining = max(0, limit_val - used_val)
                        label = "credit" if remaining == 1 else "credits"
                        pct_used = 0
                        try:
                            if limit_val > 0:
                                pct_used = int(round((used_val / float(limit_val)) * 100))
                        except Exception:
                            pct_used = 0
                        usage_nav = {
                            'text': f"{remaining} {label} left",
                            'title': f"{used_val} of {limit_val} used this month · {remaining} {label} remaining",
                            'pct': pct_used,
                        }
                    else:
                        usage_nav = { 'text': '∞', 'title': 'Unlimited this month', 'pct': 0 }
                except Exception:
                    usage_nav = None
        if email:
            try:
                import hashlib
                fb_user_match = hashlib.sha256((email or '').strip().lower().encode('utf-8', 'ignore')).hexdigest()
            except Exception:
                fb_user_match = None
        def stripe_price_url(pid: str|None):
            if not pid:
                return '#'
            key = os.getenv('STRIPE_SECRET_KEY','')
            base = 'https://dashboard.stripe.com/prices/'
            try:
                if 'sk_test' in key:
                    base = 'https://dashboard.stripe.com/test/prices/'
            except Exception:
                pass
            return base + str(pid)
        def env(name, default=''):
            return os.getenv(name, default)

        def pack_effective_price_id(c):
            """Per-collection priceId overrides STRIPE_DEFAULT_PACK_PRICE; empty => no Buy."""
            pid = (c.get('priceId') or '').strip()
            if not pid:
                pid = os.getenv('STRIPE_DEFAULT_PACK_PRICE', '').strip()
            return pid or None

        def month_key():
            """Get current month key for usage tracking"""
            return datetime.now(timezone.utc).strftime('%Y-%m')

        return {
            'fb_purchase': fb_purchase,
            'fb_user_match': fb_user_match,
            'is_admin': is_admin_email(email),
            'is_signed_in': bool(email),
            'is_pro': is_pro,
            'support_email': os.getenv('SUPPORT_EMAIL', 'support@faithsparksprintables.com'),
            'stripe_pk': STRIPE_PUBLISHABLE_KEY,
            'theme_name': theme_name,
            'theme': theme_vars,
            'logo_url': logo_url,
            'favicon_url': favicon_url,
            'site_content': site_content,
            'usage_nav': usage_nav,
            'env': env,
            'pack_effective_price_id': pack_effective_price_id,
            'stripe_price_url': stripe_price_url,
            'month_key': month_key,
            'plan_label': plan_label,
            'current_year': datetime.now().year,
            'static_v': STATIC_VERSION,

        }
    except Exception:
        return {
            'fb_purchase': fb_purchase,
            'fb_user_match': fb_user_match,
            'is_admin': False,
            'is_signed_in': False,
            'is_pro': False,
            'support_email': os.getenv('SUPPORT_EMAIL', 'support@faithsparksprintables.com'),
            'stripe_pk': None,
            'theme_name': 'teal',
            'theme': THEMES.get('teal'),
            'logo_url': url_for('static', filename='faith_sparks_logo.png'),
            'favicon_url': url_for('static', filename='favicon.ico'),
            'site_content': {},
            'usage_nav': None,
            'static_v': STATIC_VERSION,
        }

# --- Plus / Checkout ---
@app.route('/plus')
def plus_pricing():
    from faithsparks.views.billing import plus_pricing as _impl
    return _impl()

def _is_safe_next(target: str) -> bool:
    if not target:
        return False
    host_url = urlparse(request.host_url)
    test_url = urlparse(urljoin(request.host_url, target))
    return (test_url.scheme in ("http", "https") and host_url.netloc == test_url.netloc)


def resolve_price_id_local(id_or_product: str) -> str:
    """Accepts a price_... or prod_... and returns a valid price id.
    If a product id is given, tries product.default_price else the first active recurring price.
    """
    if not stripe:
        raise RuntimeError('Stripe SDK not available')
    pid = (id_or_product or '').strip()
    if not pid:
        raise ValueError('Missing price id')
    if pid.startswith('price_'):
        return pid
    if pid.startswith('prod_'):
        # Try default_price first
        try:
            prod = stripe.Product.retrieve(pid, expand=['default_price'])
            dp = prod.get('default_price')
            if isinstance(dp, dict) and dp.get('id'):
                return dp['id']
            # fallback: list active recurring prices
            prices = stripe.Price.list(product=pid, active=True, limit=10)
            if prices and prices.data:
                # prefer recurring, else first
                recurring = [p for p in prices.data if p.get('recurring')]
                target = (recurring[0] if recurring else prices.data[0])
                return target.id
        except Exception:
            pass
        # As a final fallback, let it error out
        return pid
    # Unknown format; let Stripe validate
    return pid

@app.route('/create_checkout_session', methods=['POST'])
@login_required
def create_checkout_session():
    from faithsparks.views.billing import create_checkout_session as _impl
    return _impl()

@app.route('/plus/success')
def plus_success():
    session.pop('_uc', None)  # plan likely just changed — force a fresh read
    from faithsparks.views.billing import plus_success as _impl
    return _impl()

@app.route('/billing')
@login_required
def billing_portal():
    from faithsparks.views.billing import billing_portal as _impl
    return _impl()

@app.route('/stripe/webhook', methods=['POST'])
def stripe_webhook():
    from faithsparks.views.billing import stripe_webhook as _impl
    return _impl()

@app.route('/admin/analytics')
@admin_required
def admin_analytics():
    from faithsparks.views.admin.analytics import admin_analytics as _impl
    return _impl()

# ----- Admin: Gift Plan -----
@app.route('/admin/gift', methods=['GET','POST'])
@admin_required
def admin_gift():
    from faithsparks.views.admin.gifts import admin_gift as _impl
    return _impl()

# ----- Admin: Collections CRUD -----
@app.route('/admin/collections')
@admin_required
def admin_collections():
    from faithsparks.views.admin.collections import admin_collections as _impl
    return _impl()

@app.route('/admin/collections/new', methods=['GET','POST'])
@admin_required
def admin_collections_new():
    from faithsparks.views.admin.collections import admin_collections_new as _impl
    return _impl()

@app.route('/admin/collections/<slug>', methods=['GET','POST'])
@admin_required
def admin_collections_edit(slug):
    from faithsparks.views.admin.collections import admin_collections_edit as _impl
    return _impl(slug)

@app.route('/admin/collections/<slug>/delete', methods=['POST'])
@admin_required
def admin_collections_delete(slug):
    from faithsparks.views.admin.collections import admin_collections_delete as _impl
    return _impl(slug)

@app.route('/admin/theme', methods=['GET','POST'])
@admin_required
def admin_theme():
    from faithsparks.views.admin.theme import admin_theme as _impl
    return _impl()

@app.route('/admin/theme/new', methods=['GET','POST'])
@admin_required
def admin_theme_new():
    from faithsparks.views.admin.theme import admin_theme_new as _impl
    return _impl()

@app.route('/admin/theme/<slug>', methods=['GET','POST'])
@admin_required
def admin_theme_edit(slug):
    from faithsparks.views.admin.theme import admin_theme_edit as _impl
    return _impl(slug)

@app.route('/admin/theme/<slug>/delete', methods=['POST'])
@admin_required
def admin_theme_delete(slug):
    from faithsparks.views.admin.theme import admin_theme_delete as _impl
    return _impl(slug)

def _parse_date_str(s: str):
    from datetime import date
    s = (s or '').strip()
    if not s:
        return None
    try:
        # Try YYYY-MM-DD
        parts = [int(p) for p in s.split('-')]
        if len(parts) == 3:
            return date(parts[0], parts[1], parts[2])
    except Exception:
        pass
    try:
        # Try MM-DD (assume this year)
        parts = [int(p) for p in s.split('-')]
        if len(parts) == 2:
            today = datetime.now(timezone.utc).date()
            return today.replace(month=parts[0], day=parts[1])
    except Exception:
        return None
    return None

def _is_today_in_range(start_s: str, end_s: str) -> bool:
    from datetime import date
    today = datetime.now(timezone.utc).date()
    start = _parse_date_str(start_s)
    end = _parse_date_str(end_s)
    if not start or not end:
        return False
    # Normalize years to this or adjacent year to handle wrapping (e.g., Dec -> Jan)
    s = start.replace(year=today.year)
    e = end.replace(year=today.year)
    if e < s:
        # wraps year-end: if today >= s or today <= e
        return today >= s or today <= e
    return s <= today <= e

def _parse_time_str(s: str):
    s = (s or '').strip()
    if not s:
        return None
    try:
        hh, mm = s.split(':')
        return int(hh) * 60 + int(mm)
    except Exception:
        return None

def _is_now_in_time_range(start_t: str, end_t: str) -> bool:
    """Return True if current local time is within [start,end], inclusive; supports wrap-around (e.g., 22:00–06:00)."""
    now = datetime.now().hour * 60 + datetime.now().minute
    s = _parse_time_str(start_t)
    e = _parse_time_str(end_t)
    if s is None or e is None:
        return True
    if e < s:
        return now >= s or now <= e
    return s <= now <= e

def _weekday_today_sun0() -> int:
    # Python: Monday=0..Sunday=6; convert to Sunday=0..Saturday=6
    return (datetime.now().weekday() + 1) % 7

def _rule_matches(r: dict) -> bool:
    if not r or not r.get('name'):
        return False
    if not _is_today_in_range(r.get('start',''), r.get('end','')):
        return False
    if not _is_now_in_time_range(r.get('timeStart',''), r.get('timeEnd','')):
        return False
    w = r.get('weekdays') or []
    if isinstance(w, list) and len(w):
        if _weekday_today_sun0() not in [int(x) for x in w]:
            return False
    return True

@app.route('/admin/theme/auto', methods=['POST'])
@admin_required
def admin_theme_auto():
    from faithsparks.views.admin.theme import admin_theme_auto as _impl
    return _impl()

def _save_auto_rules(rules: list[dict]):
    if not db:
        return
    db.collection('config').document('app').set({ 'autoThemes': rules }, merge=True)

@app.route('/admin/theme/auto_rules/add', methods=['POST'])
@admin_required
def admin_theme_add_rule():
    from faithsparks.views.admin.theme import admin_theme_add_rule as _impl
    return _impl()

@app.route('/admin/theme/auto_rules/update', methods=['POST'])
@admin_required
def admin_theme_update_rule():
    from faithsparks.views.admin.theme import admin_theme_update_rule as _impl
    return _impl()

@app.route('/admin/theme/auto_rules/delete', methods=['POST'])
@admin_required
def admin_theme_delete_rule():
    from faithsparks.views.admin.theme import admin_theme_delete_rule as _impl
    return _impl()

@app.route('/admin/theme/vars/<name>')
@admin_required
def admin_theme_vars(name):
    from faithsparks.views.admin.theme import admin_theme_vars as _impl
    return _impl(name)

@app.route('/admin/theme/logo', methods=['POST'])
@admin_required
def admin_theme_logo():
    from faithsparks.views.admin.theme import admin_theme_logo as _impl
    return _impl()

@app.route('/admin/theme/favicon', methods=['POST'])
@admin_required
def admin_theme_favicon():
    from faithsparks.views.admin.theme import admin_theme_favicon as _impl
    return _impl()

@app.route('/admin/content', methods=['GET','POST'])
@admin_required
def admin_content():
    from faithsparks.views.admin.content import admin_content as _impl
    return _impl()

@app.route('/admin/help')
@admin_required
def admin_help():
    from faithsparks.views.admin.help import admin_help as _impl
    return _impl()

@app.route('/admin/theme/clone_activate', methods=['POST'])
@admin_required
def admin_theme_clone_activate():
    from faithsparks.views.admin.theme import admin_theme_clone_activate as _impl
    return _impl()

@app.route('/admin/collections/<slug>/move', methods=['POST'])
@admin_required
def admin_collections_move(slug):
    from faithsparks.views.admin.collections import admin_collections_move as _impl
    return _impl(slug)

@app.route('/admin/collections/<slug>/set_order', methods=['POST'])
@admin_required
def admin_collections_set_order(slug):
    from faithsparks.views.admin.collections import admin_collections_set_order as _impl
    return _impl(slug)

@app.route('/admin/prewarm/<slug>', methods=['POST'])
@admin_required
def admin_prewarm_pack(slug):
    from faithsparks.views.admin.collections import admin_prewarm_pack as _impl
    return _impl(slug)

@app.route('/admin/prewarm/<slug>/status')
@admin_required
def admin_prewarm_status(slug):
    from faithsparks.views.admin.collections import admin_prewarm_status as _impl
    return _impl(slug)

@app.route('/packs/<path:filename>')
@admin_required
def serve_pack(filename):
    from faithsparks.views.browse import serve_pack as _impl
    return _impl(filename)

@app.route('/dl/pack/<slug>')
def dl_pack(slug):
    from faithsparks.views.browse import dl_pack as _impl
    return _impl(slug)

@app.route('/buy/pack/<slug>')
@login_required
def buy_pack(slug):
    from faithsparks.views.billing import buy_pack as _impl
    return _impl(slug)

@app.route('/buy/success/<slug>')
@login_required
def buy_success(slug):
    session.pop('_uc', None)  # entitlement likely just changed — force a fresh read
    from faithsparks.views.billing import buy_success as _impl
    return _impl(slug)

@app.route('/toggle_favorite/<filename>', methods=['POST'])
@login_required
def toggle_favorite(filename):
    from faithsparks.views.worksheets import toggle_favorite as _impl
    return _impl(filename)

@app.route('/browse')
def browse():
    from faithsparks.views.browse import browse as _impl
    return _impl()

@app.route('/browse/<slug>')
def browse_detail(slug):
    from faithsparks.views.browse import browse_detail as _impl
    return _impl(slug)


@app.route('/games')
def games():
    from faithsparks.views.games import games as _impl
    return _impl()


@app.route('/games/<slug>')
def games_detail(slug):
    from faithsparks.views.games import games_detail as _impl
    return _impl(slug)

@app.route('/games/create', methods=['GET', 'POST'])
def games_create():
    from faithsparks.views.games import games_create as _impl
    return _impl()


@app.route('/games/download/<slug>')
def dl_game(slug):
    from faithsparks.views.games import dl_game as _impl
    return _impl(slug)

@app.post('/games/words')
def games_words():
    from faithsparks.views.games import games_words as _impl
    return _impl()

@app.post("/admin/reset_credits/<uid>")
@admin_required
def admin_reset_credits(uid):
    # Fetch user document
    ref = db.collection("users").document(uid)
    snap = ref.get()
    user = snap.to_dict() if snap.exists else {}

    # Determine plan
    plan = user.get("plan", "free")  # "free", "plus_family", "plus_classroom"

    if plan == "free":
        credits = {
            "lifetime": 10,   # reset lifetime to full
            "monthly": 1,     # 1 per month
        }
    elif plan == "plus_family":
        credits = {
            "monthly": 15,
        }
    elif plan == "plus_classroom":
        credits = {
            "monthly": 100,
        }
    else:
        credits = {}  # fallback if unknown plan

    # Store credits back
    ref.set({"credits": credits}, merge=True)

    flash(f"Credits reset for {uid} ({plan})", "success")
    return redirect(url_for("admin_users"))

@app.route("/regenerate/<filename>")
@login_required
def regenerate(filename):
    from faithsparks.views.worksheets import regenerate as _impl
    return _impl(filename)


# --- Error Handlers ---
@app.errorhandler(403)
def handle_403(e):
    try:
        return render_template("403.html"), 403
    except Exception:
        return "Forbidden (CSRF)", 403


@app.errorhandler(404)
def handle_404(e):
    try:
        return render_template('404.html'), 404
    except Exception:
        return "Not found", 404

@app.errorhandler(500)
def handle_500(e):
    try:
        return render_template('500.html'), 500
    except Exception:
        return "Server error", 500

def plan_norm(p: str) -> str:
    """Normalize plan names to canonical values."""
    p = (p or "free").lower()
    if p in ("plus_family", "family", "plus"):  # treat legacy 'plus' as family
        return "plus_family"
    if p in ("plus_classroom", "classroom", "school"):
        return "plus_classroom"
    return "free"

def plan_label(p: str) -> str:
    """Get human-readable label for a plan."""
    return {
        "plus_family": "Plus Family",
        "plus_classroom": "Plus Classroom", 
        "free": "Free",
    }.get(plan_norm(p), p or "Free")
@app.route("/admin/users")
@admin_required
def admin_users():
    from faithsparks.views.admin.users import admin_users as _impl
    return _impl()

@app.route("/admin/users/<uid>/set_plan", methods=["POST"])
@admin_required 
def admin_users_set_plan(uid):
    from faithsparks.views.admin.users import admin_users_set_plan as _impl
    return _impl(uid)

@app.get("/health")
def health():
    return {"ok": True}, 200

@app.get("/api/usage")
def api_usage():
    email = session.get("user_email")
    if not db or not email:
        return jsonify({"text": "", "title": ""}), 200
    plan = _get_user_plan(email)
    m_lim, _ = _quota_for_plan(plan)
    used_life, used_m = _get_usage(email)
    if m_lim is not None:
        try:
            used_val = int(used_m)
        except Exception:
            used_val = 0
        try:
            limit_val = int(m_lim)
        except Exception:
            limit_val = 0
        remaining = max(0, limit_val - used_val)
        label = "credit" if remaining == 1 else "credits"
        data = {
            "text": f"{remaining} {label} left",
            "title": f"{used_val} of {limit_val} used this month · {remaining} {label} remaining",
        }
    else:
        data = {"text": "∞", "title": "Unlimited this month"}
    resp = jsonify(data)
    resp.headers["Cache-Control"] = "no-store, max-age=0"
    resp.headers["Pragma"] = "no-cache"
    return resp


@app.route("/admin/users/<uid>/reset_usage", methods=["POST"])
@admin_required
def admin_users_reset_usage(uid):
    from faithsparks.views.admin.users import admin_users_reset_usage as _impl
    return _impl(uid)
